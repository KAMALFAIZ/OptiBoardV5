"""Export API routes - PDF et Excel"""
from fastapi import APIRouter, Query, HTTPException
from fastapi.responses import FileResponse
from datetime import date, datetime
from typing import Optional
import os
import tempfile
import pandas as pd
from io import BytesIO

from ..database_unified import execute_app as execute_query, execute_central as execute_master_query, DWHConnectionManager
from ..sql.query_templates import (
    CHIFFRE_AFFAIRES_PAR_PERIODE,
    CHIFFRE_AFFAIRES_PAR_GAMME,
    CHIFFRE_AFFAIRES_PAR_COMMERCIAL,
    TOP_CLIENTS,
    STOCK_PAR_ARTICLE,
    STOCK_DORMANT,
    BALANCE_AGEE,
    TOP_ENCOURS_CLIENTS
)
from ..services.calculs import get_periode_dates

router = APIRouter(prefix="/api/export", tags=["Export"])

# Répertoire temporaire pour les exports
EXPORT_DIR = tempfile.gettempdir()


@router.get("/excel/ventes")
async def export_ventes_excel(
    date_debut: Optional[date] = Query(None),
    date_fin: Optional[date] = Query(None),
    periode: Optional[str] = Query("annee_courante")
):
    """
    Exporte les données de ventes en Excel.
    """
    try:
        if not date_debut or not date_fin:
            date_debut_str, date_fin_str = get_periode_dates(periode)
        else:
            date_debut_str = date_debut.strftime("%Y-%m-%d")
            date_fin_str = date_fin.strftime("%Y-%m-%d")

        # Récupérer les données
        par_periode = execute_query(CHIFFRE_AFFAIRES_PAR_PERIODE, (date_debut_str, date_fin_str))
        par_gamme = execute_query(CHIFFRE_AFFAIRES_PAR_GAMME, (date_debut_str, date_fin_str))
        par_commercial = execute_query(CHIFFRE_AFFAIRES_PAR_COMMERCIAL, (date_debut_str, date_fin_str))
        top_clients = execute_query(TOP_CLIENTS, (date_debut_str, date_fin_str))

        # Créer le fichier Excel avec plusieurs feuilles
        filename = f"ventes_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        filepath = os.path.join(EXPORT_DIR, filename)

        with pd.ExcelWriter(filepath, engine='openpyxl') as writer:
            # Feuille 1: CA par période
            if par_periode:
                df_periode = pd.DataFrame(par_periode)
                df_periode.to_excel(writer, sheet_name='CA par Période', index=False)

            # Feuille 2: CA par gamme
            if par_gamme:
                df_gamme = pd.DataFrame(par_gamme)
                df_gamme.to_excel(writer, sheet_name='CA par Gamme', index=False)

            # Feuille 3: CA par commercial
            if par_commercial:
                df_commercial = pd.DataFrame(par_commercial)
                df_commercial.to_excel(writer, sheet_name='CA par Commercial', index=False)

            # Feuille 4: Top clients
            if top_clients:
                df_clients = pd.DataFrame(top_clients)
                df_clients.to_excel(writer, sheet_name='Top Clients', index=False)

        return FileResponse(
            filepath,
            media_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            filename=filename
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/excel/stocks")
async def export_stocks_excel():
    """
    Exporte les données de stocks en Excel.
    """
    try:
        # Récupérer les données
        stock_articles = execute_query(STOCK_PAR_ARTICLE)
        stock_dormant = execute_query(STOCK_DORMANT)

        filename = f"stocks_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        filepath = os.path.join(EXPORT_DIR, filename)

        with pd.ExcelWriter(filepath, engine='openpyxl') as writer:
            if stock_articles:
                df_articles = pd.DataFrame(stock_articles)
                df_articles.to_excel(writer, sheet_name='Stock par Article', index=False)

            if stock_dormant:
                df_dormant = pd.DataFrame(stock_dormant)
                df_dormant.to_excel(writer, sheet_name='Stock Dormant', index=False)

        return FileResponse(
            filepath,
            media_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            filename=filename
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/excel/recouvrement")
async def export_recouvrement_excel():
    """
    Exporte les données de recouvrement en Excel.
    """
    try:
        balance = execute_query(BALANCE_AGEE)
        top_encours = execute_query(TOP_ENCOURS_CLIENTS)

        filename = f"recouvrement_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        filepath = os.path.join(EXPORT_DIR, filename)

        with pd.ExcelWriter(filepath, engine='openpyxl') as writer:
            if balance:
                df_balance = pd.DataFrame(balance)
                df_balance.to_excel(writer, sheet_name='Balance Âgée', index=False)

            if top_encours:
                df_encours = pd.DataFrame(top_encours)
                df_encours.to_excel(writer, sheet_name='Top Encours', index=False)

        return FileResponse(
            filepath,
            media_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            filename=filename
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/excel/complet")
async def export_rapport_complet_excel(
    date_debut: Optional[date] = Query(None),
    date_fin: Optional[date] = Query(None),
    periode: Optional[str] = Query("annee_courante")
):
    """
    Exporte un rapport complet en Excel (toutes les données).
    """
    try:
        if not date_debut or not date_fin:
            date_debut_str, date_fin_str = get_periode_dates(periode)
        else:
            date_debut_str = date_debut.strftime("%Y-%m-%d")
            date_fin_str = date_fin.strftime("%Y-%m-%d")

        # Toutes les données
        par_periode = execute_query(CHIFFRE_AFFAIRES_PAR_PERIODE, (date_debut_str, date_fin_str))
        par_gamme = execute_query(CHIFFRE_AFFAIRES_PAR_GAMME, (date_debut_str, date_fin_str))
        par_commercial = execute_query(CHIFFRE_AFFAIRES_PAR_COMMERCIAL, (date_debut_str, date_fin_str))
        top_clients = execute_query(TOP_CLIENTS, (date_debut_str, date_fin_str))
        stock_articles = execute_query(STOCK_PAR_ARTICLE)
        stock_dormant = execute_query(STOCK_DORMANT)
        balance = execute_query(BALANCE_AGEE)

        filename = f"rapport_complet_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
        filepath = os.path.join(EXPORT_DIR, filename)

        with pd.ExcelWriter(filepath, engine='openpyxl') as writer:
            # Synthèse
            synthese_data = {
                "Indicateur": [
                    "CA Total HT",
                    "Nombre de Clients",
                    "Nombre d'Articles en Stock",
                    "Valeur Stock Dormant",
                    "Encours Clients Total"
                ],
                "Valeur": [
                    sum(row.get('CA_HT', 0) or 0 for row in par_periode),
                    len(top_clients),
                    len(stock_articles),
                    sum(row.get('Valeur_Stock', 0) or 0 for row in stock_dormant),
                    sum(row.get('Solde_Cloture', 0) or 0 for row in balance)
                ]
            }
            df_synthese = pd.DataFrame(synthese_data)
            df_synthese.to_excel(writer, sheet_name='Synthèse', index=False)

            # Autres feuilles
            if par_periode:
                pd.DataFrame(par_periode).to_excel(writer, sheet_name='CA Mensuel', index=False)
            if par_gamme:
                pd.DataFrame(par_gamme).to_excel(writer, sheet_name='CA par Gamme', index=False)
            if par_commercial:
                pd.DataFrame(par_commercial).to_excel(writer, sheet_name='CA par Commercial', index=False)
            if top_clients:
                pd.DataFrame(top_clients).to_excel(writer, sheet_name='Top Clients', index=False)
            if stock_articles:
                pd.DataFrame(stock_articles).to_excel(writer, sheet_name='Stocks', index=False)
            if stock_dormant:
                pd.DataFrame(stock_dormant).to_excel(writer, sheet_name='Stock Dormant', index=False)
            if balance:
                pd.DataFrame(balance).to_excel(writer, sheet_name='Balance Âgée', index=False)

        return FileResponse(
            filepath,
            media_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            filename=filename
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/pdf/dashboard")
async def export_dashboard_pdf(
    date_debut: Optional[date] = Query(None),
    date_fin: Optional[date] = Query(None),
    periode: Optional[str] = Query("annee_courante")
):
    """
    Génère un PDF du dashboard.
    """
    try:
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4, landscape
        from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import cm

        if not date_debut or not date_fin:
            date_debut_str, date_fin_str = get_periode_dates(periode)
        else:
            date_debut_str = date_debut.strftime("%Y-%m-%d")
            date_fin_str = date_fin.strftime("%Y-%m-%d")

        # Récupérer les données
        par_periode = execute_query(CHIFFRE_AFFAIRES_PAR_PERIODE, (date_debut_str, date_fin_str))
        balance = execute_query(BALANCE_AGEE)

        # Calculs
        ca_total = sum(row.get('CA_HT', 0) or 0 for row in par_periode)
        encours_total = sum(row.get('Solde_Cloture', 0) or 0 for row in balance)

        filename = f"dashboard_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
        filepath = os.path.join(EXPORT_DIR, filename)

        doc = SimpleDocTemplate(filepath, pagesize=A4)
        styles = getSampleStyleSheet()
        elements = []

        # Titre
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=18,
            spaceAfter=30,
            alignment=1  # Centré
        )
        elements.append(Paragraph("Reporting Commercial - KAsoft", title_style))
        elements.append(Paragraph(f"Période: {date_debut_str} au {date_fin_str}", styles['Normal']))
        elements.append(Spacer(1, 20))

        # KPIs
        kpi_data = [
            ["Indicateur", "Valeur"],
            ["Chiffre d'Affaires HT", f"{ca_total:,.2f} MAD"],
            ["Encours Clients", f"{encours_total:,.2f} MAD"],
            ["Nombre de Clients", str(len(balance))],
        ]

        kpi_table = Table(kpi_data, colWidths=[8*cm, 6*cm])
        kpi_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#1e40af')),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
            ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -1), colors.HexColor('#f8fafc')),
            ('GRID', (0, 0), (-1, -1), 1, colors.HexColor('#e2e8f0'))
        ]))
        elements.append(kpi_table)
        elements.append(Spacer(1, 30))

        # CA Mensuel
        elements.append(Paragraph("Evolution du CA Mensuel", styles['Heading2']))
        elements.append(Spacer(1, 10))

        if par_periode:
            ca_mensuel_data = [["Période", "CA HT", "Marge Brute", "Nb Clients"]]
            for row in par_periode[:12]:
                ca_mensuel_data.append([
                    f"{row['Annee']}-{str(row['Mois']).zfill(2)}",
                    f"{row.get('CA_HT', 0) or 0:,.2f}",
                    f"{(row.get('CA_HT', 0) or 0) - (row.get('Cout_Total', 0) or 0):,.2f}",
                    str(row.get('Nb_Clients', 0) or 0)
                ])

            ca_table = Table(ca_mensuel_data, colWidths=[4*cm, 4*cm, 4*cm, 3*cm])
            ca_table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#059669')),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.white),
                ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
                ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                ('FONTSIZE', (0, 0), (-1, -1), 9),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.HexColor('#d1d5db'))
            ]))
            elements.append(ca_table)

        # Footer
        elements.append(Spacer(1, 30))
        elements.append(Paragraph(
            f"Généré le {datetime.now().strftime('%d/%m/%Y à %H:%M')}",
            styles['Normal']
        ))

        doc.build(elements)

        return FileResponse(
            filepath,
            media_type='application/pdf',
            filename=filename
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/pptx/dashboard")
async def export_dashboard_pptx(
    date_debut: Optional[date] = Query(None),
    date_fin: Optional[date] = Query(None),
    periode: Optional[str] = Query("annee_courante")
):
    """Génère un PowerPoint du tableau de bord exécutif."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        from io import BytesIO

        BLUE      = RGBColor(0x25, 0x63, 0xEB)
        DARK_BLUE = RGBColor(0x1E, 0x40, 0xAF)
        GREEN     = RGBColor(0x05, 0x96, 0x69)
        ORANGE    = RGBColor(0xD9, 0x77, 0x06)
        WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
        GRAY_LIGHT = RGBColor(0xF3, 0xF4, 0xF6)

        if not date_debut or not date_fin:
            date_debut_str, date_fin_str = get_periode_dates(periode)
        else:
            date_debut_str = date_debut.strftime("%Y-%m-%d")
            date_fin_str   = date_fin.strftime("%Y-%m-%d")

        par_periode    = execute_query(CHIFFRE_AFFAIRES_PAR_PERIODE,    (date_debut_str, date_fin_str))
        par_commercial = execute_query(CHIFFRE_AFFAIRES_PAR_COMMERCIAL, (date_debut_str, date_fin_str))
        top_clients    = execute_query(TOP_CLIENTS,                     (date_debut_str, date_fin_str))
        balance        = execute_query(BALANCE_AGEE)

        ca_total      = sum(r.get('CA_HT', 0) or 0 for r in par_periode)
        encours_total = sum(r.get('Solde_Cloture', 0) or 0 for r in balance)
        nb_clients    = len(top_clients)

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        def add_text(slide, text, left, top, width, height, size=12, bold=False, color=None, align=PP_ALIGN.LEFT):
            txb = slide.shapes.add_textbox(left, top, width, height)
            tf  = txb.text_frame
            tf.word_wrap = True
            p   = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = bold
            if color:
                run.font.color.rgb = color

        # ── Slide 1 : Titre ─────────────────────────────────────────────────
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = DARK_BLUE

        add_text(slide, "Tableau de Bord Exécutif", Inches(1), Inches(2.2), Inches(11.33), Inches(1.3),
                 size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_text(slide, f"Période : {date_debut_str}  →  {date_fin_str}",
                 Inches(1), Inches(3.7), Inches(11.33), Inches(0.6),
                 size=16, color=RGBColor(0xBF,0xDB,0xFE), align=PP_ALIGN.CENTER)
        add_text(slide, f"Généré le {datetime.now().strftime('%d/%m/%Y à %H:%M')}  —  OptiBoard",
                 Inches(1), Inches(4.5), Inches(11.33), Inches(0.5),
                 size=11, color=RGBColor(0x93,0xC5,0xFD), align=PP_ALIGN.CENTER)

        # ── Slide 2 : KPIs ──────────────────────────────────────────────────
        slide = prs.slides.add_slide(blank)
        add_text(slide, "Indicateurs Clés", Inches(0.4), Inches(0.2), Inches(8), Inches(0.5),
                 size=18, bold=True, color=DARK_BLUE)
        add_text(slide, f"{date_debut_str}  →  {date_fin_str}",
                 Inches(0.4), Inches(0.7), Inches(8), Inches(0.35), size=10, color=RGBColor(0x6B,0x72,0x80))

        def kpi_card(slide, label, value, left, top, color):
            from pptx.util import Emu
            from pptx.enum.shapes import PP_PLACEHOLDER
            w, h = Inches(3.8), Inches(2.2)
            # MSO_AUTO_SHAPE_TYPE.RECTANGLE = 1
            box = slide.shapes.add_shape(1, left, top, w, h)
            box.fill.solid(); box.fill.fore_color.rgb = color
            box.line.fill.background()
            txb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.25), w - Inches(0.4), Inches(0.5))
            p = txb.text_frame.paragraphs[0]
            r = p.add_run(); r.text = label
            r.font.size = Pt(13); r.font.color.rgb = RGBColor(0xFF,0xFF,0xFF); r.font.bold = False
            txb2 = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.9), w - Inches(0.4), Inches(0.9))
            p2 = txb2.text_frame.paragraphs[0]
            r2 = p2.add_run(); r2.text = value
            r2.font.size = Pt(26); r2.font.bold = True; r2.font.color.rgb = WHITE

        kpi_card(slide, "Chiffre d'Affaires HT",  f"{ca_total:,.0f} MAD",       Inches(0.4),  Inches(1.3), BLUE)
        kpi_card(slide, "Encours Clients",          f"{encours_total:,.0f} MAD", Inches(4.75), Inches(1.3), GREEN)
        kpi_card(slide, "Nombre de Clients actifs", str(nb_clients),             Inches(9.1),  Inches(1.3), ORANGE)

        # ── Slide 3 : CA Mensuel ─────────────────────────────────────────────
        if par_periode:
            slide = prs.slides.add_slide(blank)
            add_text(slide, "Évolution du Chiffre d'Affaires Mensuel",
                     Inches(0.4), Inches(0.2), Inches(12), Inches(0.5), size=18, bold=True, color=DARK_BLUE)

            headers = ["Période", "CA HT", "Coût", "Marge Brute", "Nb Clients"]
            rows_data = []
            for r in par_periode[:20]:
                ca   = r.get('CA_HT', 0) or 0
                cout = r.get('Cout_Total', 0) or 0
                rows_data.append([
                    f"{r['Annee']}-{str(r['Mois']).zfill(2)}",
                    f"{ca:,.2f}",
                    f"{cout:,.2f}",
                    f"{ca-cout:,.2f}",
                    str(r.get('Nb_Clients', 0) or 0)
                ])

            nc = len(headers)
            nr = len(rows_data) + 1
            table = slide.shapes.add_table(nr, nc, Inches(0.4), Inches(0.9), Inches(12.53), Inches(6.3)).table
            col_widths = [Inches(1.8), Inches(2.5), Inches(2.5), Inches(2.5), Inches(1.5)]
            for ci in range(nc):
                table.columns[ci].width = col_widths[ci] if ci < len(col_widths) else Inches(2)

            for ci, h in enumerate(headers):
                cell = table.cell(0, ci)
                cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
                p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = h; r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = WHITE

            for ri, row_vals in enumerate(rows_data):
                for ci, val in enumerate(row_vals):
                    cell = table.cell(ri+1, ci)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = GRAY_LIGHT if ri % 2 == 0 else WHITE
                    p = cell.text_frame.paragraphs[0]
                    p.alignment = PP_ALIGN.RIGHT if ci > 0 else PP_ALIGN.LEFT
                    r = p.add_run(); r.text = val; r.font.size = Pt(8)

        # ── Slide 4 : Top Clients ────────────────────────────────────────────
        if top_clients:
            slide = prs.slides.add_slide(blank)
            add_text(slide, "Top Clients par Chiffre d'Affaires",
                     Inches(0.4), Inches(0.2), Inches(12), Inches(0.5), size=18, bold=True, color=DARK_BLUE)

            top20 = top_clients[:20]
            col_keys = list(top20[0].keys()) if top20 else []
            nc = min(len(col_keys), 6)
            col_keys = col_keys[:nc]
            nr = len(top20) + 1

            table = slide.shapes.add_table(nr, nc, Inches(0.4), Inches(0.9), Inches(12.53), Inches(6.3)).table
            cw = Inches(12.53) // nc
            for ci in range(nc):
                table.columns[ci].width = cw

            for ci, key in enumerate(col_keys):
                cell = table.cell(0, ci)
                cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
                p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = str(key); r.font.size = Pt(9); r.font.bold = True; r.font.color.rgb = WHITE

            for ri, row_data in enumerate(top20):
                for ci, key in enumerate(col_keys):
                    val = row_data.get(key, "")
                    cell = table.cell(ri+1, ci)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = GRAY_LIGHT if ri % 2 == 0 else WHITE
                    p = cell.text_frame.paragraphs[0]
                    is_num = isinstance(val, (int, float)) and not isinstance(val, bool)
                    p.alignment = PP_ALIGN.RIGHT if is_num else PP_ALIGN.LEFT
                    r = p.add_run()
                    r.text = f"{val:,.2f}" if is_num else str(val if val is not None else "")
                    r.font.size = Pt(8)

        filename = f"dashboard_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        filepath = os.path.join(EXPORT_DIR, filename)

        output = BytesIO()
        prs.save(output)
        output.seek(0)

        with open(filepath, 'wb') as f:
            f.write(output.getvalue())

        return FileResponse(
            filepath,
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            filename=filename
        )

    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/csv/{table}")
async def export_csv(
    table: str,
    date_debut: Optional[date] = Query(None),
    date_fin: Optional[date] = Query(None)
):
    """
    Exporte une table spécifique en CSV.
    """
    try:
        tables_map = {
            "ventes": CHIFFRE_AFFAIRES_PAR_PERIODE,
            "gammes": CHIFFRE_AFFAIRES_PAR_GAMME,
            "commerciaux": CHIFFRE_AFFAIRES_PAR_COMMERCIAL,
            "stocks": STOCK_PAR_ARTICLE,
            "dormant": STOCK_DORMANT,
            "balance": BALANCE_AGEE
        }

        if table not in tables_map:
            raise HTTPException(status_code=400, detail="Table non valide")

        query = tables_map[table]

        if date_debut and date_fin:
            date_debut_str = date_debut.strftime("%Y-%m-%d")
            date_fin_str = date_fin.strftime("%Y-%m-%d")
            if "?" in query:
                data = execute_query(query, (date_debut_str, date_fin_str))
            else:
                data = execute_query(query)
        else:
            date_debut_str, date_fin_str = get_periode_dates("annee_courante")
            if "?" in query:
                data = execute_query(query, (date_debut_str, date_fin_str))
            else:
                data = execute_query(query)

        if not data:
            raise HTTPException(status_code=404, detail="Aucune donnée")

        df = pd.DataFrame(data)

        filename = f"{table}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.csv"
        filepath = os.path.join(EXPORT_DIR, filename)

        df.to_csv(filepath, index=False, encoding='utf-8-sig')

        return FileResponse(
            filepath,
            media_type='text/csv',
            filename=filename
        )

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# ==================== FONCTION POUR LE SCHEDULER ====================


def _execute_template_query(query: str, date_debut: str, date_fin: str, societe: str = None) -> list:
    """
    Execute une requete de datasource template sur la bonne base DWH.
    Utilise inject_params pour les parametres et DWHConnectionManager pour la connexion.
    """
    from ..services.parameter_resolver import inject_params

    # Construire le contexte des parametres
    context = {
        'dateDebut': date_debut,
        'dateFin': date_fin,
        'societe': societe  # None = toutes les societes
    }

    # Injecter les parametres dans la requete
    final_query = inject_params(query, context)

    # Trouver le DWH actif (toujours dans MASTER)
    dwh_list = execute_master_query(
        "SELECT TOP 1 code FROM APP_DWH WHERE actif = 1 ORDER BY id",
        use_cache=True
    )

    if dwh_list:
        return DWHConnectionManager.execute_dwh_query(dwh_list[0]['code'], final_query, use_cache=False)
    else:
        return execute_query(final_query, use_cache=False)


async def generate_report_file(report_type: str, report_id: int, export_format: str, filters: dict = None) -> dict:
    """
    Genere un fichier de rapport pour l'envoi par email.
    Utilisee par le scheduler.

    Args:
        report_type: pivot, gridview, dashboard, export
        report_id: ID du rapport ou nom pour les exports standards
        export_format: excel, pdf, csv
        filters: Filtres optionnels (dates, etc.)

    Returns:
        dict avec success, file_path et report_name
    """
    try:
        from ..database_unified import execute_app as execute_query
        import json

        date_debut_str, date_fin_str = get_periode_dates(filters.get('period', 'annee_courante') if filters else 'annee_courante')

        if filters and filters.get('date_debut'):
            date_debut_str = filters['date_debut']
        if filters and filters.get('date_fin'):
            date_fin_str = filters['date_fin']

        report_name = "Rapport"
        filepath = None

        # Export standard (ventes, stocks, etc.)
        if report_type == 'export':
            if report_id == 'ventes' or report_id == 1:
                report_name = "Export Ventes"
                par_periode = execute_query(CHIFFRE_AFFAIRES_PAR_PERIODE, (date_debut_str, date_fin_str))
                par_gamme = execute_query(CHIFFRE_AFFAIRES_PAR_GAMME, (date_debut_str, date_fin_str))
                par_commercial = execute_query(CHIFFRE_AFFAIRES_PAR_COMMERCIAL, (date_debut_str, date_fin_str))
                top_clients = execute_query(TOP_CLIENTS, (date_debut_str, date_fin_str))

                filename = f"ventes_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
                filepath = os.path.join(EXPORT_DIR, filename)

                with pd.ExcelWriter(filepath, engine='openpyxl') as writer:
                    if par_periode:
                        pd.DataFrame(par_periode).to_excel(writer, sheet_name='CA par Periode', index=False)
                    if par_gamme:
                        pd.DataFrame(par_gamme).to_excel(writer, sheet_name='CA par Gamme', index=False)
                    if par_commercial:
                        pd.DataFrame(par_commercial).to_excel(writer, sheet_name='CA par Commercial', index=False)
                    if top_clients:
                        pd.DataFrame(top_clients).to_excel(writer, sheet_name='Top Clients', index=False)

            elif report_id == 'stocks' or report_id == 2:
                report_name = "Export Stocks"
                stock_articles = execute_query(STOCK_PAR_ARTICLE)
                stock_dormant = execute_query(STOCK_DORMANT)

                filename = f"stocks_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
                filepath = os.path.join(EXPORT_DIR, filename)

                with pd.ExcelWriter(filepath, engine='openpyxl') as writer:
                    if stock_articles:
                        pd.DataFrame(stock_articles).to_excel(writer, sheet_name='Stock par Article', index=False)
                    if stock_dormant:
                        pd.DataFrame(stock_dormant).to_excel(writer, sheet_name='Stock Dormant', index=False)

            elif report_id == 'recouvrement' or report_id == 3:
                report_name = "Export Recouvrement"
                balance = execute_query(BALANCE_AGEE)
                top_encours = execute_query(TOP_ENCOURS_CLIENTS)

                filename = f"recouvrement_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
                filepath = os.path.join(EXPORT_DIR, filename)

                with pd.ExcelWriter(filepath, engine='openpyxl') as writer:
                    if balance:
                        pd.DataFrame(balance).to_excel(writer, sheet_name='Balance Agee', index=False)
                    if top_encours:
                        pd.DataFrame(top_encours).to_excel(writer, sheet_name='Top Encours', index=False)

            elif report_id == 'complet' or report_id == 4:
                report_name = "Rapport Complet"
                par_periode = execute_query(CHIFFRE_AFFAIRES_PAR_PERIODE, (date_debut_str, date_fin_str))
                par_gamme = execute_query(CHIFFRE_AFFAIRES_PAR_GAMME, (date_debut_str, date_fin_str))
                par_commercial = execute_query(CHIFFRE_AFFAIRES_PAR_COMMERCIAL, (date_debut_str, date_fin_str))
                top_clients = execute_query(TOP_CLIENTS, (date_debut_str, date_fin_str))
                stock_articles = execute_query(STOCK_PAR_ARTICLE)
                stock_dormant = execute_query(STOCK_DORMANT)
                balance = execute_query(BALANCE_AGEE)

                filename = f"rapport_complet_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
                filepath = os.path.join(EXPORT_DIR, filename)

                with pd.ExcelWriter(filepath, engine='openpyxl') as writer:
                    if par_periode:
                        pd.DataFrame(par_periode).to_excel(writer, sheet_name='CA Mensuel', index=False)
                    if par_gamme:
                        pd.DataFrame(par_gamme).to_excel(writer, sheet_name='CA par Gamme', index=False)
                    if par_commercial:
                        pd.DataFrame(par_commercial).to_excel(writer, sheet_name='CA par Commercial', index=False)
                    if top_clients:
                        pd.DataFrame(top_clients).to_excel(writer, sheet_name='Top Clients', index=False)
                    if stock_articles:
                        pd.DataFrame(stock_articles).to_excel(writer, sheet_name='Stocks', index=False)
                    if stock_dormant:
                        pd.DataFrame(stock_dormant).to_excel(writer, sheet_name='Stock Dormant', index=False)
                    if balance:
                        pd.DataFrame(balance).to_excel(writer, sheet_name='Balance Agee', index=False)

        # Pivot
        elif report_type in ('pivot', 'pivot-v2'):
            pivot = execute_query("SELECT * FROM APP_Pivots WHERE id = ?", (report_id,), use_cache=False)
            if not pivot:
                return {"success": False, "error": "Pivot non trouvé"}

            pivot = pivot[0]
            report_name = pivot['nom']

            # Recuperer les donnees - chercher dans APP_DataSources puis APP_DataSources_Templates
            query = None
            ds_id = pivot.get('data_source_id')
            ds_code = pivot.get('data_source_code')

            if ds_id:
                ds = execute_query("SELECT query_template FROM APP_DataSources WHERE id = ?", (ds_id,), use_cache=False)
                if ds and ds[0].get('query_template'):
                    query = ds[0]['query_template']

            if not query and ds_code:
                ds = execute_master_query("SELECT query_template FROM APP_DataSources_Templates WHERE code = ?", (ds_code,), use_cache=False)
                if ds and ds[0].get('query_template'):
                    query = ds[0]['query_template']

            if query:
                data = _execute_template_query(query, date_debut_str, date_fin_str)

                filename = f"pivot_{report_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
                filepath = os.path.join(EXPORT_DIR, filename)

                df = pd.DataFrame(data) if data else pd.DataFrame()
                df.to_excel(filepath, sheet_name=report_name[:31], index=False)

        # GridView
        elif report_type == 'gridview':
            gridview = execute_query("SELECT * FROM APP_GridViews WHERE id = ?", (report_id,), use_cache=False)
            if not gridview:
                return {"success": False, "error": "GridView non trouvé"}

            gridview = gridview[0]
            report_name = gridview['nom']

            # Recuperer les donnees - chercher dans APP_DataSources puis APP_DataSources_Templates
            query = None
            ds_code = gridview.get('data_source_code')
            ds_id = gridview.get('data_source_id')

            # 1. Essayer par data_source_id dans APP_DataSources
            if ds_id:
                ds = execute_query("SELECT query_template FROM APP_DataSources WHERE id = ?", (ds_id,), use_cache=False)
                if ds and ds[0].get('query_template'):
                    query = ds[0]['query_template']

            # 2. Sinon essayer par data_source_code dans APP_DataSources_Templates
            if not query and ds_code:
                ds = execute_master_query("SELECT query_template FROM APP_DataSources_Templates WHERE code = ?", (ds_code,), use_cache=False)
                if ds and ds[0].get('query_template'):
                    query = ds[0]['query_template']

            if query:
                data = _execute_template_query(query, date_debut_str, date_fin_str)

                filename = f"gridview_{report_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"
                filepath = os.path.join(EXPORT_DIR, filename)

                df = pd.DataFrame(data) if data else pd.DataFrame()
                df.to_excel(filepath, sheet_name=report_name[:31], index=False)

        # Dashboard - generer PDF
        elif report_type == 'dashboard':
            dashboard = execute_query("SELECT * FROM APP_Dashboards WHERE id = ?", (report_id,), use_cache=False)
            if not dashboard:
                return {"success": False, "error": "Dashboard non trouve"}

            report_name = dashboard[0]['nom']

            # Generer un PDF simple du dashboard
            from reportlab.lib import colors
            from reportlab.lib.pagesizes import A4
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
            from reportlab.lib.styles import getSampleStyleSheet

            filename = f"dashboard_{report_id}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pdf"
            filepath = os.path.join(EXPORT_DIR, filename)

            doc = SimpleDocTemplate(filepath, pagesize=A4)
            styles = getSampleStyleSheet()
            elements = []

            elements.append(Paragraph(f"Dashboard: {report_name}", styles['Heading1']))
            elements.append(Spacer(1, 20))
            elements.append(Paragraph(f"Periode: {date_debut_str} au {date_fin_str}", styles['Normal']))
            elements.append(Spacer(1, 10))
            elements.append(Paragraph(f"Genere le: {datetime.now().strftime('%d/%m/%Y %H:%M')}", styles['Normal']))

            doc.build(elements)

        if not filepath or not os.path.exists(filepath):
            return {"success": False, "error": "Fichier non genere"}

        return {
            "success": True,
            "file_path": filepath,
            "report_name": report_name
        }

    except Exception as e:
        return {"success": False, "error": str(e)}
