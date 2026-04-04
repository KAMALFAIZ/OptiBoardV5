"""
Générateur de documents (PowerPoint & Excel) basé sur des templates.
L'utilisateur choisit un template, remplit un formulaire, l'IA enrichit le contenu.
"""
from fastapi import APIRouter, HTTPException
from fastapi.responses import StreamingResponse
from pydantic import BaseModel
from typing import Dict, Any
import json, io, logging

logger = logging.getLogger(__name__)
router = APIRouter(prefix="/api/ai/presentation", tags=["AI Presentation Builder"])


# ── Modèles ───────────────────────────────────────────────────────────────────

class GenerateRequest(BaseModel):
    template_id: str
    form_data: Dict[str, str] = {}
    doc_type: str = "pptx"   # "pptx" | "excel"


# ── Helpers ───────────────────────────────────────────────────────────────────

def _split_bullets(text: str, defaults: list, max_count: int = 4) -> list:
    """Transforme un champ texte libre en liste de bullets."""
    if not text or not text.strip():
        return defaults[:max_count]
    lines = [l.strip() for l in text.replace(';', '\n').replace(',', '\n').splitlines() if l.strip()]
    return (lines[:max_count] if lines else defaults[:max_count])


# ── Prompts IA par template ───────────────────────────────────────────────────

def _build_pptx_prompt(template_name: str, form_data: Dict) -> str:
    fields = "\n".join(f"- {k}: {v}" for k, v in form_data.items() if v and k != 'style')
    style = form_data.get('style', 'formal')
    return f"""Tu es un expert en présentations professionnelles en français.
Génère le contenu JSON pour une présentation PowerPoint "{template_name}".

Informations fournies par l'utilisateur:
{fields}

Réponds UNIQUEMENT avec ce JSON valide (aucun texte avant ou après):
{{
  "title": "titre principal",
  "subtitle": "sous-titre court",
  "style": "{style}",
  "audience": "public cible déduit des informations",
  "slides": [
    {{"type": "title"}},
    {{"type": "agenda", "items": ["Section 1", "Section 2", "Section 3", "Section 4"]}},
    {{"type": "content", "title": "Titre slide 1", "bullets": ["Point précis 1", "Point précis 2", "Point précis 3"], "notes": ""}},
    {{"type": "section", "title": "Deuxième partie"}},
    {{"type": "content", "title": "Titre slide 2", "bullets": ["Point A", "Point B", "Point C"], "notes": ""}},
    {{"type": "content", "title": "Titre slide 3", "bullets": ["Point X", "Point Y", "Point Z"], "notes": ""}},
    {{"type": "closing", "message": "Message de conclusion adapté"}}
  ]
}}

Génère 6 à 8 slides. Utilise les informations fournies. Contenu professionnel en français."""


def _build_excel_prompt(template_name: str, form_data: Dict) -> str:
    fields = "\n".join(f"- {k}: {v}" for k, v in form_data.items() if v)
    return f"""Tu es un expert Excel professionnel en français.
Génère le contenu JSON pour un fichier Excel "{template_name}".

Informations:
{fields}

Réponds UNIQUEMENT avec ce JSON valide:
{{
  "title": "Titre du fichier",
  "sheet_name": "Nom feuille",
  "columns": [
    {{"name": "Colonne 1", "type": "text", "width": 20}},
    {{"name": "Colonne 2", "type": "number", "width": 15}},
    {{"name": "Colonne 3", "type": "currency", "width": 15}},
    {{"name": "Colonne 4", "type": "percent", "width": 12}}
  ],
  "sample_rows": [
    ["Valeur A", 100, 5000, 0.15],
    ["Valeur B", 200, 8500, 0.22],
    ["Valeur C", 150, 6200, 0.18],
    ["Valeur D", 300, 12000, 0.25],
    ["Valeur E", 250, 9800, 0.20]
  ],
  "has_total_row": true,
  "has_chart": true,
  "chart_type": "bar"
}}

Génère 5-8 lignes de données réalistes et cohérentes. Types disponibles: text, number, currency, date, percent."""


# ── Fallbacks sans IA ─────────────────────────────────────────────────────────

def _fallback_pptx(template_id: str, form_data: Dict) -> Dict:
    company  = form_data.get('company', 'Notre Entreprise')
    period   = form_data.get('period', '')
    year     = form_data.get('year', '2025')
    style    = form_data.get('style', 'formal')

    configs = {
        'rapport_mensuel': {
            'title': company,
            'subtitle': f'Rapport {period}' if period else 'Rapport Mensuel',
            'slides': [
                {'type': 'title'},
                {'type': 'agenda', 'items': ['Résultats du mois', 'Analyse des performances', 'Points d\'attention', 'Actions à venir']},
                {'type': 'content', 'title': 'Résultats du mois',
                 'bullets': _split_bullets(form_data.get('highlights'), ['Chiffre d\'affaires en progression', 'Objectifs partiellement atteints', 'Nouveaux clients acquis', 'Marge maintenue'])},
                {'type': 'content', 'title': 'Analyse des performances',
                 'bullets': ['Comparaison vs objectifs', 'Évolution vs période précédente', 'Analyse par segment', 'Facteurs clés de succès']},
                {'type': 'section', 'title': 'Actions à venir'},
                {'type': 'content', 'title': 'Plan d\'actions',
                 'bullets': _split_bullets(form_data.get('next_actions'), ['Renforcer les efforts commerciaux', 'Suivi des opportunités en cours', 'Réunion équipe vendredi', 'Rapport direction'])},
                {'type': 'closing', 'message': 'Merci pour votre attention'},
            ]
        },
        'bilan_annuel': {
            'title': company,
            'subtitle': f'Bilan {year}',
            'slides': [
                {'type': 'title'},
                {'type': 'agenda', 'items': ['Points forts', 'Défis rencontrés', 'Résultats clés', 'Perspectives']},
                {'type': 'content', 'title': f'Points forts {year}',
                 'bullets': _split_bullets(form_data.get('highlights'), ['Croissance du chiffre d\'affaires', 'Élargissement de la clientèle', 'Amélioration de la satisfaction client', 'Innovation produit'])},
                {'type': 'content', 'title': 'Défis & enseignements',
                 'bullets': _split_bullets(form_data.get('challenges'), ['Contexte marché tendu', 'Pression sur les marges', 'Recrutement difficile', 'Concurrence accrue'])},
                {'type': 'section', 'title': 'Perspectives'},
                {'type': 'content', 'title': 'Objectifs & perspectives',
                 'bullets': _split_bullets(form_data.get('outlook'), ['Objectif croissance +15%', 'Nouveaux marchés cibles', 'Investissements prévus', 'Plan de recrutement'])},
                {'type': 'closing', 'message': f"Merci · Cap sur {int(year)+1 if year.isdigit() else 'l avenir'}"},
            ]
        },
        'pitch_client': {
            'title': form_data.get('client', 'Notre Client'),
            'subtitle': form_data.get('project', 'Proposition Commerciale'),
            'slides': [
                {'type': 'title'},
                {'type': 'agenda', 'items': ['Qui sommes-nous ?', 'Votre besoin', 'Notre solution', 'Bénéfices & ROI', 'Prochaines étapes']},
                {'type': 'content', 'title': f'Qui sommes-nous — {company}',
                 'bullets': ['Leader dans notre domaine', 'Plus de 10 ans d\'expertise', 'Références clients solides', 'Équipe dédiée à votre succès']},
                {'type': 'content', 'title': 'Votre besoin identifié',
                 'bullets': _split_bullets(form_data.get('value_prop'), ['Optimisation des processus', 'Réduction des coûts', 'Gain de temps significatif', 'Amélioration de la performance'])},
                {'type': 'section', 'title': 'Notre Solution'},
                {'type': 'content', 'title': f'{form_data.get("project", "Notre Solution")}',
                 'bullets': ['Approche sur mesure', 'Déploiement rapide', 'Support dédié inclus', 'Résultats mesurables']},
                {'type': 'content', 'title': 'Prochaines étapes',
                 'bullets': [form_data.get('next_steps', 'Réunion de lancement'), 'Validation du périmètre', 'Signature du contrat', 'Démarrage du projet']},
                {'type': 'closing', 'message': 'Construisons ensemble votre succès'},
            ]
        },
        'presentation_produit': {
            'title': form_data.get('product', 'Notre Produit'),
            'subtitle': f'Par {company}',
            'slides': [
                {'type': 'title'},
                {'type': 'agenda', 'items': ['Le problème', 'Notre solution', 'Fonctionnalités', 'Bénéfices', 'Tarifs & démarrage']},
                {'type': 'content', 'title': 'Le problème que nous résolvons',
                 'bullets': ['Processus trop longs et coûteux', 'Manque de visibilité en temps réel', 'Outils disparates non intégrés', f'Impact sur {form_data.get("target", "votre activité")}']},
                {'type': 'content', 'title': f'{form_data.get("product", "Notre Solution")} — en bref',
                 'bullets': _split_bullets(form_data.get('features'), ['Interface intuitive', 'Déploiement en 48h', 'Intégration avec vos outils', 'Support 7j/7'])},
                {'type': 'section', 'title': 'Bénéfices & ROI'},
                {'type': 'content', 'title': 'Pourquoi choisir notre solution',
                 'bullets': [f'Cible : {form_data.get("target", "Toutes entreprises")}', 'ROI dès le premier mois', 'Satisfaction client > 95%', form_data.get('pricing', 'Tarif sur devis')]},
                {'type': 'closing', 'message': 'Démarrez aujourd\'hui — essai gratuit 30 jours'},
            ]
        },
        'revue_performance': {
            'title': form_data.get('team', 'Revue de Performance'),
            'subtitle': f'{company} · {period}',
            'slides': [
                {'type': 'title'},
                {'type': 'agenda', 'items': ['KPIs & résultats', 'Points forts', 'Axes d\'amélioration', 'Plan d\'action']},
                {'type': 'content', 'title': 'KPIs & résultats',
                 'bullets': _split_bullets(form_data.get('kpis'), ['Objectif 1 : atteint à 110%', 'Objectif 2 : en progression', 'Objectif 3 : à renforcer', 'NPS : +42'])},
                {'type': 'content', 'title': 'Points forts de la période',
                 'bullets': ['Cohésion d\'équipe renforcée', 'Délais respectés', 'Qualité de livraison', 'Initiative et autonomie']},
                {'type': 'section', 'title': 'Plan d\'amélioration'},
                {'type': 'content', 'title': 'Axes d\'amélioration',
                 'bullets': _split_bullets(form_data.get('improvements'), ['Communication inter-équipes', 'Montée en compétences', 'Optimisation des processus', 'Suivi des indicateurs'])},
                {'type': 'closing', 'message': 'Ensemble, continuons à progresser'},
            ]
        },
        'plan_strategique': {
            'title': company,
            'subtitle': f'Plan Stratégique {form_data.get("horizon", "2025-2027")}',
            'slides': [
                {'type': 'title'},
                {'type': 'agenda', 'items': ['Vision', 'Axes stratégiques', 'Feuille de route', 'Indicateurs de succès']},
                {'type': 'content', 'title': 'Notre vision',
                 'bullets': _split_bullets(form_data.get('vision'), ['Leader reconnu sur notre marché', 'Croissance durable et responsable', 'Innovation au cœur de nos activités', 'Satisfaction client au centre'])},
                {'type': 'content', 'title': 'Axes stratégiques',
                 'bullets': _split_bullets(form_data.get('axes'), ['Développement commercial', 'Excellence opérationnelle', 'Innovation & digital', 'Développement des talents'])},
                {'type': 'section', 'title': 'Feuille de route'},
                {'type': 'content', 'title': 'Feuille de route & priorités',
                 'bullets': ['Court terme (6 mois) : consolidation', 'Moyen terme (1 an) : expansion', 'Long terme (3 ans) : leadership', 'Revue trimestrielle des objectifs']},
                {'type': 'content', 'title': 'Indicateurs de succès',
                 'bullets': _split_bullets(form_data.get('kpis'), ['CA cible défini', 'Part de marché visée', 'NPS > 70', 'Taux de rétention > 90%'])},
                {'type': 'closing', 'message': 'Ensemble vers notre vision'},
            ]
        },
    }

    ctx = configs.get(template_id, configs['rapport_mensuel'])
    ctx['style'] = style
    ctx['audience'] = form_data.get('audience', '')
    return ctx


def _fallback_excel(template_id: str, form_data: Dict) -> Dict:
    company = form_data.get('company', 'Mon Entreprise')
    period  = form_data.get('period', '2025')

    configs = {
        'budget_previsionnel': {
            'title': f'Budget Prévisionnel {form_data.get("year", "2025")} — {company}',
            'sheet_name': 'Budget',
            'columns': [
                {'name': 'Catégorie', 'type': 'text', 'width': 25},
                {'name': 'Budget N-1', 'type': 'currency', 'width': 16},
                {'name': 'Budget N', 'type': 'currency', 'width': 16},
                {'name': 'Évolution', 'type': 'percent', 'width': 12},
                {'name': 'Commentaire', 'type': 'text', 'width': 30},
            ],
            'sample_rows': [
                ['Ressources Humaines', 320000, 350000, 0.094, 'Recrutement 3 postes'],
                ['Marketing & Communication', 85000, 95000, 0.118, 'Campagne digital'],
                ['Informatique & Digital', 60000, 80000, 0.333, 'Migration Cloud'],
                ['Frais Généraux', 45000, 47000, 0.044, 'Stable'],
                ['Formation', 22000, 28000, 0.273, 'Montée en compétences'],
                ['Commercial', 110000, 125000, 0.136, 'Ouverture région'],
            ],
            'has_total_row': True, 'has_chart': True, 'chart_type': 'bar'
        },
        'suivi_commercial': {
            'title': f'Suivi Commercial {period} — {company}',
            'sheet_name': 'Suivi',
            'columns': [
                {'name': 'Commercial', 'type': 'text', 'width': 20},
                {'name': 'Objectif CA', 'type': 'currency', 'width': 16},
                {'name': 'CA Réalisé', 'type': 'currency', 'width': 16},
                {'name': 'Atteinte', 'type': 'percent', 'width': 12},
                {'name': 'Nb Clients', 'type': 'number', 'width': 12},
                {'name': 'Nb Devis', 'type': 'number', 'width': 10},
            ],
            'sample_rows': [
                ['Jean Martin', 150000, 162000, 1.08, 18, 42],
                ['Sophie Durand', 120000, 108000, 0.90, 14, 31],
                ['Pierre Leblanc', 180000, 195000, 1.08, 22, 55],
                ['Marie Petit', 130000, 141000, 1.085, 16, 38],
                ['Thomas Bernard', 160000, 148000, 0.925, 19, 44],
            ],
            'has_total_row': True, 'has_chart': True, 'chart_type': 'bar'
        },
        'analyse_clients': {
            'title': f'Analyse Clients — {company}',
            'sheet_name': 'Clients',
            'columns': [
                {'name': 'Client', 'type': 'text', 'width': 25},
                {'name': 'Segment', 'type': 'text', 'width': 15},
                {'name': 'CA Annuel', 'type': 'currency', 'width': 16},
                {'name': 'Évolution', 'type': 'percent', 'width': 12},
                {'name': 'Nb Commandes', 'type': 'number', 'width': 14},
                {'name': 'Panier Moyen', 'type': 'currency', 'width': 14},
            ],
            'sample_rows': [
                ['Groupe Alpha', 'Grand Compte', 285000, 0.12, 38, 7500],
                ['Beta Industries', 'ETI', 142000, 0.08, 24, 5917],
                ['Gamma SARL', 'PME', 68000, 0.22, 18, 3778],
                ['Delta Corp', 'Grand Compte', 198000, -0.05, 31, 6387],
                ['Epsilon SA', 'ETI', 95000, 0.15, 20, 4750],
                ['Zeta Group', 'PME', 52000, 0.31, 12, 4333],
            ],
            'has_total_row': True, 'has_chart': True, 'chart_type': 'bar'
        },
        'planning_projet': {
            'title': f'{form_data.get("project", "Projet")} — Planning',
            'sheet_name': 'Planning',
            'columns': [
                {'name': 'Phase / Tâche', 'type': 'text', 'width': 30},
                {'name': 'Responsable', 'type': 'text', 'width': 18},
                {'name': 'Début', 'type': 'date', 'width': 14},
                {'name': 'Fin', 'type': 'date', 'width': 14},
                {'name': 'Avancement', 'type': 'percent', 'width': 12},
                {'name': 'Statut', 'type': 'text', 'width': 15},
            ],
            'sample_rows': [
                ['Phase 1 — Cadrage', form_data.get('team', 'Chef de projet'), '2025-01-06', '2025-01-17', 1.0, 'Terminé'],
                ['Phase 2 — Conception', form_data.get('team', 'Équipe design'), '2025-01-20', '2025-02-07', 0.8, 'En cours'],
                ['Phase 3 — Développement', form_data.get('team', 'Équipe tech'), '2025-02-10', '2025-03-21', 0.4, 'En cours'],
                ['Phase 4 — Tests', form_data.get('team', 'QA'), '2025-03-24', '2025-04-04', 0.0, 'À venir'],
                ['Phase 5 — Déploiement', form_data.get('team', 'DevOps'), '2025-04-07', '2025-04-11', 0.0, 'À venir'],
                ['Phase 6 — Formation', form_data.get('team', 'Formateur'), '2025-04-14', '2025-04-25', 0.0, 'À venir'],
            ],
            'has_total_row': False, 'has_chart': False, 'chart_type': 'bar'
        },
    }
    return configs.get(template_id, configs['suivi_commercial'])


# ── Route principale ──────────────────────────────────────────────────────────

@router.post("/generate")
async def generate_document(req: GenerateRequest):
    """Génère un document à partir d'un template et des données du formulaire."""
    from ..services.ai_provider import get_ai_provider, AIMessage

    provider = get_ai_provider()
    template_names = {
        'rapport_mensuel': 'Rapport Mensuel',    'bilan_annuel': 'Bilan Annuel',
        'pitch_client': 'Pitch Client',           'presentation_produit': 'Présentation Produit',
        'revue_performance': 'Revue de Performance', 'plan_strategique': 'Plan Stratégique',
        'budget_previsionnel': 'Budget Prévisionnel', 'suivi_commercial': 'Suivi Commercial',
        'analyse_clients': 'Analyse Clients',     'planning_projet': 'Planning Projet',
    }
    template_name = template_names.get(req.template_id, req.template_id)

    if req.doc_type == 'pptx':
        # Essayer l'IA, sinon fallback
        ctx = None
        if provider:
            try:
                prompt = _build_pptx_prompt(template_name, req.form_data)
                messages = [
                    AIMessage("system", "Tu génères uniquement du JSON valide, rien d'autre."),
                    AIMessage("user", prompt)
                ]
                raw = await provider.chat(messages)
                raw = raw.strip()
                if "```json" in raw:
                    raw = raw.split("```json")[1].split("```")[0].strip()
                elif "```" in raw:
                    raw = raw.split("```")[1].split("```")[0].strip()
                ctx = json.loads(raw)
                # S'assurer que le style du formulaire est respecté
                ctx['style'] = req.form_data.get('style', ctx.get('style', 'formal'))
            except Exception as e:
                logger.warning(f"AI generation failed, using fallback: {e}")
                ctx = None

        if ctx is None:
            ctx = _fallback_pptx(req.template_id, req.form_data)

        output, filename = _build_pptx(ctx)
        return StreamingResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'}
        )

    else:  # excel
        ctx = None
        if provider:
            try:
                prompt = _build_excel_prompt(template_name, req.form_data)
                messages = [
                    AIMessage("system", "Tu génères uniquement du JSON valide, rien d'autre."),
                    AIMessage("user", prompt)
                ]
                raw = await provider.chat(messages)
                raw = raw.strip()
                if "```json" in raw:
                    raw = raw.split("```json")[1].split("```")[0].strip()
                elif "```" in raw:
                    raw = raw.split("```")[1].split("```")[0].strip()
                ctx = json.loads(raw)
            except Exception as e:
                logger.warning(f"AI Excel generation failed, using fallback: {e}")
                ctx = None

        if ctx is None:
            ctx = _fallback_excel(req.template_id, req.form_data)

        output, filename = _build_excel(ctx)
        return StreamingResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition": f'attachment; filename="{filename}"'}
        )


# ── Générateur PPTX ───────────────────────────────────────────────────────────

def _build_pptx(ctx: dict):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from datetime import datetime

    style = ctx.get("style", "formal")
    palettes = {
        "formal":  {"bg": RGBColor(0x1E,0x40,0xAF), "accent": RGBColor(0x25,0x63,0xEB), "light": RGBColor(0xDB,0xEA,0xFE), "text": RGBColor(0x1E,0x29,0x4E)},
        "dynamic": {"bg": RGBColor(0xEA,0x58,0x0C), "accent": RGBColor(0xF9,0x73,0x16), "light": RGBColor(0xFF,0xED,0xD5), "text": RGBColor(0x43,0x18,0x07)},
        "minimal": {"bg": RGBColor(0x18,0x18,0x18), "accent": RGBColor(0x52,0x52,0x52), "light": RGBColor(0xF5,0xF5,0xF5), "text": RGBColor(0x18,0x18,0x18)},
    }
    pal   = palettes.get(style, palettes["formal"])
    WHITE = RGBColor(0xFF,0xFF,0xFF)
    GRAY  = RGBColor(0x6B,0x72,0x80)
    LGRAY = RGBColor(0xF9,0xFA,0xFB)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    title_main   = ctx.get("title", "Présentation")
    subtitle_txt = ctx.get("subtitle", "")
    slides_data  = ctx.get("slides", [])

    def add_rect(slide, l, t, w, h, color):
        sh = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background()

    def add_text(slide, text, l, t, w, h, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
        txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = txb.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = str(text)
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color

    for sd in slides_data:
        stype = sd.get("type", "content")
        slide = prs.slides.add_slide(blank)

        if stype == "title":
            add_rect(slide, 0, 0, 13.33, 7.5, pal["bg"])
            add_rect(slide, 0, 0, 0.5, 7.5, pal["accent"])
            add_text(slide, title_main, 1.2, 2.0, 11, 1.5, 40, bold=True)
            if subtitle_txt:
                add_text(slide, subtitle_txt, 1.2, 3.7, 10, 0.7, 18, color=RGBColor(0xBF,0xDB,0xFE))
            meta = [ctx.get("audience",""), datetime.now().strftime("%B %Y")]
            add_text(slide, "  ·  ".join(m for m in meta if m), 1.2, 6.2, 10, 0.5, 11,
                     color=RGBColor(0x93,0xC5,0xFD), italic=True)
            add_text(slide, "OptiBoard", 11.5, 6.9, 1.6, 0.4, 9,
                     color=RGBColor(0x60,0x82,0xBD), italic=True, align=PP_ALIGN.RIGHT)

        elif stype == "agenda":
            add_rect(slide, 0, 0, 13.33, 1.1, pal["bg"])
            add_text(slide, "Sommaire", 0.5, 0.15, 12, 0.8, 22, bold=True)
            items = sd.get("items", [])
            for i, item in enumerate(items):
                col, row = i % 2, i // 2
                l, t = 0.8 + col * 6.3, 1.4 + row * 0.9
                add_rect(slide, l, t + 0.06, 0.42, 0.44, pal["accent"])
                add_text(slide, str(i+1), l + 0.05, t, 0.32, 0.56, 14, bold=True, align=PP_ALIGN.CENTER)
                add_text(slide, item, l + 0.55, t + 0.08, 5.4, 0.5, 13, color=pal["text"])

        elif stype == "section":
            add_rect(slide, 0, 0, 13.33, 7.5, pal["light"])
            add_rect(slide, 0, 2.9, 13.33, 1.7, pal["bg"])
            add_rect(slide, 0, 0, 0.3, 7.5, pal["accent"])
            add_text(slide, sd.get("title", ""), 0.8, 3.0, 11.5, 1.5, 30, bold=True, align=PP_ALIGN.CENTER)

        elif stype == "content":
            add_rect(slide, 0, 0, 13.33, 1.05, pal["bg"])
            add_rect(slide, 0, 1.05, 13.33, 0.06, pal["accent"])
            add_text(slide, sd.get("title", ""), 0.4, 0.12, 12.5, 0.85, 22, bold=True)
            add_rect(slide, 0, 1.11, 13.33, 6.39, LGRAY)
            for i, bullet in enumerate(sd.get("bullets", [])[:6]):
                t = 1.55 + i * 0.76
                add_rect(slide, 0.4, t + 0.18, 0.16, 0.16, pal["accent"])
                add_text(slide, bullet, 0.73, t, 12.0, 0.7, 14, color=pal["text"])
            if sd.get("notes"):
                slide.notes_slide.notes_text_frame.text = sd["notes"]
            add_text(slide, title_main, 0.4, 7.1, 9, 0.32, 8, color=GRAY, italic=True)
            add_text(slide, datetime.now().strftime("%d/%m/%Y"), 11.5, 7.1, 1.6, 0.32, 8,
                     color=GRAY, italic=True, align=PP_ALIGN.RIGHT)

        elif stype == "closing":
            add_rect(slide, 0, 0, 13.33, 7.5, pal["bg"])
            add_rect(slide, 0, 0, 0.5, 7.5, pal["accent"])
            msg = sd.get("message", "Merci pour votre attention")
            add_text(slide, msg, 1.2, 2.6, 10.5, 1.4, 32, bold=True)
            add_text(slide, "Questions & Échanges", 1.2, 4.2, 8, 0.6, 16, color=RGBColor(0xBF,0xDB,0xFE))
            add_text(slide, "OptiBoard", 11.5, 6.9, 1.6, 0.4, 9,
                     color=RGBColor(0x60,0x82,0xBD), italic=True, align=PP_ALIGN.RIGHT)

    buf = io.BytesIO(); prs.save(buf); buf.seek(0)
    safe = "".join(c for c in title_main if c.isalnum() or c in " _-")[:40].strip().replace(" ", "_")
    return buf, f"{safe or 'presentation'}.pptx"


# ── Générateur Excel ──────────────────────────────────────────────────────────

def _build_excel(ctx: dict):
    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    except ImportError:
        raise HTTPException(status_code=500, detail="openpyxl non installé.")

    from datetime import datetime

    wb  = openpyxl.Workbook()
    ws  = wb.active
    ws.title = ctx.get("sheet_name", "Données")[:31]

    title_txt = ctx.get("title", "Rapport")
    columns   = ctx.get("columns", [])
    rows      = ctx.get("sample_rows", [])
    has_total = ctx.get("has_total_row", False)
    has_chart = ctx.get("has_chart", False)

    def cf(bold=False, color="000000", size=11, italic=False):
        return Font(bold=bold, color=color, size=size, italic=italic, name="Calibri")

    def fill(hex_c):
        return PatternFill("solid", fgColor=hex_c)

    def border():
        s = Side(style="thin", color="D1D5DB")
        return Border(left=s, right=s, top=s, bottom=s)

    # Titre
    ws.merge_cells(f"A1:{openpyxl.utils.get_column_letter(max(len(columns),1))}1")
    ws["A1"] = title_txt
    ws["A1"].font      = cf(bold=True, color="FFFFFF", size=15)
    ws["A1"].fill      = fill("1E3A8A")
    ws["A1"].alignment = Alignment(horizontal="center", vertical="center")
    ws.row_dimensions[1].height = 34

    ws.merge_cells(f"A2:{openpyxl.utils.get_column_letter(max(len(columns),1))}2")
    ws["A2"] = f"Généré par OptiBoard · {datetime.now().strftime('%d/%m/%Y à %H:%M')}"
    ws["A2"].font      = cf(color="6B7280", size=9, italic=True)
    ws["A2"].fill      = fill("EFF6FF")
    ws["A2"].alignment = Alignment(horizontal="center")
    ws.row_dimensions[2].height = 16
    ws.row_dimensions[3].height = 6

    HR = 4
    col_types = []
    for ci, col in enumerate(columns, 1):
        c = ws.cell(row=HR, column=ci, value=col.get("name", f"Col {ci}"))
        c.font      = cf(bold=True, color="FFFFFF", size=11)
        c.fill      = fill("2563EB")
        c.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        c.border    = border()
        ws.column_dimensions[c.column_letter].width = col.get("width", 15)
        col_types.append(col.get("type", "text"))
    ws.row_dimensions[HR].height = 26

    FMT = {"number": "#,##0.00", "currency": '#,##0.00 "€"', "percent": "0.00%",
           "date": "DD/MM/YYYY", "text": "@"}

    for ri, row_data in enumerate(rows):
        er  = HR + 1 + ri
        bg  = "F9FAFB" if ri % 2 == 0 else "FFFFFF"
        for ci, val in enumerate(row_data[:len(columns)], 1):
            c = ws.cell(row=er, column=ci, value=val)
            ct = col_types[ci-1] if ci-1 < len(col_types) else "text"
            if ct in FMT:
                c.number_format = FMT[ct]
            c.fill      = fill(bg)
            c.border    = border()
            c.alignment = Alignment(
                horizontal="right" if ct in ("number","currency","percent") else "left",
                vertical="center")
        ws.row_dimensions[er].height = 19

    if has_total and rows:
        tr = HR + 1 + len(rows)
        ws.cell(row=tr, column=1, value="TOTAL").font = cf(bold=True, color="FFFFFF")
        ws.cell(row=tr, column=1).fill = fill("1E3A8A")
        ws.cell(row=tr, column=1).border = border()
        for ci, ct in enumerate(col_types, 1):
            if ct in ("number","currency","percent"):
                cl = ws.cell(row=HR+1, column=ci).column_letter
                c  = ws.cell(row=tr, column=ci)
                c.value          = f"=SUM({cl}{HR+1}:{cl}{HR+len(rows)})"
                c.font           = cf(bold=True, color="FFFFFF")
                c.fill           = fill("1E3A8A")
                c.border         = border()
                c.alignment      = Alignment(horizontal="right")
                c.number_format  = FMT.get(ct, "")
            else:
                c = ws.cell(row=tr, column=ci)
                c.fill = fill("1E3A8A"); c.border = border()
        ws.row_dimensions[tr].height = 21

    if has_chart and rows and len(columns) >= 2:
        try:
            from openpyxl.chart import BarChart, Reference
            chart = BarChart()
            chart.type   = "col"; chart.style = 10
            chart.title  = title_txt; chart.width = 18; chart.height = 10
            data_ref = Reference(ws, min_col=2, min_row=HR, max_col=min(len(columns),4), max_row=HR+len(rows))
            cats_ref = Reference(ws, min_col=1, min_row=HR+1, max_row=HR+len(rows))
            chart.add_data(data_ref, titles_from_data=True)
            chart.set_categories(cats_ref)
            chart_col = openpyxl.utils.get_column_letter(len(columns)+2)
            ws.add_chart(chart, f"{chart_col}{HR}")
        except Exception as e:
            logger.warning(f"Chart error: {e}")

    ws.freeze_panes = ws.cell(row=HR+1, column=1)
    if columns:
        ws.auto_filter.ref = f"A{HR}:{openpyxl.utils.get_column_letter(len(columns))}{HR+len(rows)}"

    buf = io.BytesIO(); wb.save(buf); buf.seek(0)
    safe = "".join(c for c in title_txt if c.isalnum() or c in " _-")[:40].strip().replace(" ", "_")
    return buf, f"{safe or 'rapport'}.xlsx"
