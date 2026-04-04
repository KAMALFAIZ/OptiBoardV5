"""
Pivot Builder V2 - API Router
==============================
Full-stack rewrite du systeme pivot avec:
- CRUD complet des configurations
- Moteur d'execution avec sous-totaux, comparaison N/N-1
- Drill-down sur cellule
- Preferences utilisateur
- Export Excel/CSV/PDF
"""

import io
import json
import logging
import re
import time
import hashlib
from collections import defaultdict
from datetime import datetime, date
from decimal import Decimal
from typing import Dict, Any, List, Optional, Tuple

from fastapi import APIRouter, HTTPException, Header, Query
from fastapi.responses import StreamingResponse
from pydantic import BaseModel

from ..database_unified import execute_central as execute_query, central_cursor as get_db_cursor, DWHConnectionManager
from ..services.datasource_resolver import resolve_datasource
from ..services.parameter_resolver import inject_params, extract_parameters_from_query

logger = logging.getLogger("PivotV2")

router = APIRouter(prefix="/api/v2/pivots", tags=["Pivot V2"])


def _generate_entity_code(entity_type: str, nom: str) -> str:
    """Genere un code unique pour une entite (ex: PV_ventes_par_region_a3f2)"""
    slug = re.sub(r'[^a-z0-9]+', '_', nom.lower().strip())[:40].strip('_')
    suffix = hashlib.md5(f"{nom}{time.time()}".encode()).hexdigest()[:4]
    return f"{entity_type}_{slug}_{suffix}"


# =============================================================================
# PYDANTIC SCHEMAS
# =============================================================================

class PivotAxisField(BaseModel):
    field: str
    label: Optional[str] = None
    order: Optional[int] = None
    type: Optional[str] = None  # text, number, date
    date_grouping: Optional[str] = None  # jour, semaine, mois, mois_annee, trimestre, etc.
    numeric_grouping: Optional[Dict[str, Any]] = None  # {"type":"interval","step":100} ou {"type":"ranges","ranges":[...]}
    text_grouping: Optional[Dict[str, Any]] = None  # {"type":"first_letter"} ou {"type":"custom","groups":{...}}

class PivotValueField(BaseModel):
    field: str
    aggregation: str = "SUM"
    label: Optional[str] = None
    format: Optional[str] = "number"  # number, currency, percent, date, text
    decimals: Optional[int] = 2
    percentOf: Optional[str] = None  # row, column, total
    summary_aggregation: Optional[str] = None  # fonction alternative pour les totaux
    show_in_totals: Optional[bool] = True  # afficher dans les totaux

class PivotFilterField(BaseModel):
    field: str
    type: Optional[str] = "select"
    default: Optional[Any] = None
    values: Optional[List[Any]] = None

class FormattingRule(BaseModel):
    field: str
    type: str  # heatmap, data_bars, icons, thresholds, negative_red
    config: Optional[Dict[str, Any]] = None

class PivotCreateRequest(BaseModel):
    nom: str
    description: Optional[str] = None
    data_source_id: Optional[int] = None
    data_source_code: Optional[str] = None
    drilldown_data_source_code: Optional[str] = None
    rows_config: Optional[List[PivotAxisField]] = []
    columns_config: Optional[List[PivotAxisField]] = []
    filters_config: Optional[List[PivotFilterField]] = []
    values_config: Optional[List[PivotValueField]] = []
    show_grand_totals: bool = True
    show_subtotals: bool = True
    show_row_percent: bool = False
    show_col_percent: bool = False
    show_total_percent: bool = False
    grand_total_position: Optional[str] = "bottom"  # "top" ou "bottom"
    subtotal_position: Optional[str] = "bottom"  # "top" ou "bottom"
    show_summary_row: bool = False
    summary_functions: Optional[List[str]] = None  # ["AVG", "MEDIAN", "VAR"]
    comparison_mode: Optional[str] = None
    formatting_rules: Optional[List[FormattingRule]] = []
    source_params: Optional[List[Dict[str, Any]]] = []
    window_calculations: Optional[List[Dict[str, Any]]] = []
    is_public: bool = False
    application: Optional[str] = None
    created_by: Optional[int] = None

class PivotUpdateRequest(BaseModel):
    nom: Optional[str] = None
    description: Optional[str] = None
    data_source_id: Optional[int] = None
    data_source_code: Optional[str] = None
    drilldown_data_source_code: Optional[str] = None
    rows_config: Optional[List[PivotAxisField]] = None
    columns_config: Optional[List[PivotAxisField]] = None
    filters_config: Optional[List[PivotFilterField]] = None
    values_config: Optional[List[PivotValueField]] = None
    show_grand_totals: Optional[bool] = None
    show_subtotals: Optional[bool] = None
    show_row_percent: Optional[bool] = None
    show_col_percent: Optional[bool] = None
    show_total_percent: Optional[bool] = None
    grand_total_position: Optional[str] = None
    subtotal_position: Optional[str] = None
    show_summary_row: Optional[bool] = None
    summary_functions: Optional[List[str]] = None
    comparison_mode: Optional[str] = None
    formatting_rules: Optional[List[FormattingRule]] = None
    source_params: Optional[List[Dict[str, Any]]] = None
    window_calculations: Optional[List[Dict[str, Any]]] = None
    is_public: Optional[bool] = None
    application: Optional[str] = None

class PivotExecuteRequest(BaseModel):
    context: Optional[Dict[str, Any]] = {}
    raw: Optional[bool] = False

class PivotDrilldownRequest(BaseModel):
    rowValues: Dict[str, Any]
    columnValue: Optional[str] = None
    valueField: str
    context: Optional[Dict[str, Any]] = {}
    page: int = 1
    pageSize: int = 50
    sortField: Optional[str] = None
    sortDirection: Optional[str] = "asc"

class PivotUserPrefsRequest(BaseModel):
    custom_config: Optional[Dict[str, Any]] = None
    ui_state: Optional[Dict[str, Any]] = None


# =============================================================================
# INIT TABLES
# =============================================================================

def init_pivot_v2_tables():
    """Cree les tables si elles n'existent pas"""
    try:
        with get_db_cursor() as cursor:
            # APP_Pivots_V2
            cursor.execute("""
                IF NOT EXISTS (SELECT * FROM INFORMATION_SCHEMA.TABLES WHERE TABLE_NAME = 'APP_Pivots_V2')
                BEGIN
                    CREATE TABLE APP_Pivots_V2 (
                        id INT IDENTITY(1,1) PRIMARY KEY,
                        nom NVARCHAR(200) NOT NULL,
                        description NVARCHAR(500),
                        data_source_id INT,
                        data_source_code VARCHAR(100),
                        rows_config NVARCHAR(MAX),
                        columns_config NVARCHAR(MAX),
                        filters_config NVARCHAR(MAX),
                        values_config NVARCHAR(MAX),
                        show_grand_totals BIT DEFAULT 1,
                        show_subtotals BIT DEFAULT 1,
                        show_row_percent BIT DEFAULT 0,
                        show_col_percent BIT DEFAULT 0,
                        show_total_percent BIT DEFAULT 0,
                        comparison_mode VARCHAR(20) NULL,
                        formatting_rules NVARCHAR(MAX),
                        source_params NVARCHAR(MAX),
                        is_public BIT DEFAULT 0,
                        created_by INT,
                        created_at DATETIME DEFAULT GETDATE(),
                        updated_at DATETIME DEFAULT GETDATE()
                    )
                END
            """)

            # APP_Pivot_User_Prefs
            cursor.execute("""
                IF NOT EXISTS (SELECT * FROM INFORMATION_SCHEMA.TABLES WHERE TABLE_NAME = 'APP_Pivot_User_Prefs')
                BEGIN
                    CREATE TABLE APP_Pivot_User_Prefs (
                        id INT IDENTITY(1,1) PRIMARY KEY,
                        pivot_id INT NOT NULL,
                        user_id INT NOT NULL,
                        custom_config NVARCHAR(MAX),
                        ui_state NVARCHAR(MAX),
                        updated_at DATETIME DEFAULT GETDATE(),
                        CONSTRAINT FK_PivotPrefsV2_Pivot FOREIGN KEY (pivot_id)
                            REFERENCES APP_Pivots_V2(id) ON DELETE CASCADE,
                        CONSTRAINT UQ_PivotPrefsV2 UNIQUE (pivot_id, user_id)
                    )
                END
            """)
            # Ajouter les nouvelles colonnes si manquantes
            new_columns = [
                ("code", "VARCHAR(100)"),  # Code unique pour publication Master
                ("grand_total_position", "VARCHAR(10) DEFAULT 'bottom'"),
                ("subtotal_position", "VARCHAR(10) DEFAULT 'bottom'"),
                ("show_summary_row", "BIT DEFAULT 0"),
                ("summary_functions", "NVARCHAR(MAX)"),
                ("window_calculations", "NVARCHAR(MAX)"),
                ("application", "NVARCHAR(100)"),
                ("drilldown_data_source_code", "VARCHAR(100)"),
            ]
            for col_name, col_def in new_columns:
                try:
                    cursor.execute(f"""
                        IF NOT EXISTS (SELECT 1 FROM INFORMATION_SCHEMA.COLUMNS
                                       WHERE TABLE_NAME='APP_Pivots_V2' AND COLUMN_NAME='{col_name}')
                        ALTER TABLE APP_Pivots_V2 ADD {col_name} {col_def}
                    """)
                except Exception:
                    pass

        logger.info("Pivot V2 tables initialized")
    except Exception as e:
        logger.error(f"Error initializing Pivot V2 tables: {e}")


# =============================================================================
# HELPERS
# =============================================================================

def _parse_json(value: Any, default=None):
    """Parse JSON string ou retourne default"""
    if value is None:
        return default if default is not None else []
    if isinstance(value, (list, dict)):
        return value
    try:
        return json.loads(value)
    except (json.JSONDecodeError, TypeError):
        return default if default is not None else []


def _to_json(value: Any) -> Optional[str]:
    """Convertit en JSON string"""
    if value is None:
        return None
    if isinstance(value, str):
        return value
    return json.dumps(value, default=str, ensure_ascii=False)


def _normalize_row(row: Dict[str, Any]) -> Dict[str, Any]:
    """Normalise les types Decimal en float"""
    return {
        k: float(v) if isinstance(v, Decimal) else v
        for k, v in row.items()
    }


def _resolve_and_execute(
    pivot_config: Dict[str, Any],
    context: Dict[str, Any],
    dwh_code: Optional[str] = None,
    limit: Optional[int] = None,
    return_debug: bool = False
):
    """Resout la datasource et execute la requete"""
    ds_code = pivot_config.get("data_source_code")
    ds_id = pivot_config.get("data_source_id")

    # Resoudre la datasource
    if ds_code:
        ds = resolve_datasource(ds_code, dwh_code)
    elif ds_id:
        ds = resolve_datasource(ds_id, dwh_code)
    else:
        raise ValueError("Aucune source de donnees configuree pour ce pivot")

    # Injecter les parametres
    query = inject_params(ds.query_template, context)

    # Retirer ORDER BY pour usage en sous-requete
    query_clean = re.sub(r'\s+ORDER\s+BY\s+[\s\S]+$', '', query, flags=re.IGNORECASE)

    if limit:
        # Wrapper avec TOP
        query_clean = f"SELECT TOP {limit} * FROM ({query_clean}) AS __src__"

    # Executer sur le bon DWH
    effective_dwh = dwh_code
    if not effective_dwh and ds.origin.value == "template":
        try:
            dwh_list = execute_query(
                "SELECT TOP 1 code FROM APP_DWH WHERE actif = 1 ORDER BY id",
                use_cache=True
            )
            if dwh_list:
                effective_dwh = dwh_list[0]['code']
        except Exception:
            pass

    # ── DEBUG INFO ──────────────────────────────────────────────────────────
    debug_info = {
        "datasource_code": ds_code or str(ds_id),
        "datasource_origin": ds.origin.value,
        "effective_dwh": effective_dwh,
        "context_passed": context,
        "query_injected": query_clean,
    }
    logger.info(f"[PIVOT DEBUG] datasource={debug_info['datasource_code']} dwh={effective_dwh} context={context}")
    logger.info(f"[PIVOT DEBUG] query_injected=\n{query_clean}")
    # ────────────────────────────────────────────────────────────────────────

    if effective_dwh and ds.origin.value == "template":
        data = DWHConnectionManager.execute_dwh_query(effective_dwh, query_clean, use_cache=False)
    else:
        data = execute_query(query_clean, use_cache=False)

    rows = [_normalize_row(row) for row in data]

    if return_debug:
        return rows, debug_info
    return rows


# --- Regroupement temporel ---
_MONTH_NAMES_FR = {
    1: 'Janvier', 2: 'Fevrier', 3: 'Mars', 4: 'Avril',
    5: 'Mai', 6: 'Juin', 7: 'Juillet', 8: 'Aout',
    9: 'Septembre', 10: 'Octobre', 11: 'Novembre', 12: 'Decembre'
}

_MONTH_NAMES_FR_REVERSE = {v: k for k, v in _MONTH_NAMES_FR.items()}


def _date_grouping_sql_expr(field_name: str, grouping: str) -> str:
    """Retourne l'expression SQL Server equivalente au regroupement Python"""
    f = f"[{field_name}]"
    if grouping == "jour":
        return f"CONVERT(VARCHAR(10), {f}, 23)"  # YYYY-MM-DD
    elif grouping == "semaine":
        return f"CAST(YEAR({f}) AS VARCHAR) + '-S' + RIGHT('0' + CAST(DATEPART(ISO_WEEK, {f}) AS VARCHAR), 2)"
    elif grouping == "mois":
        # Retourne le numero du mois - on comparera cote Python
        return f"MONTH({f})"
    elif grouping == "mois_annee":
        # Retourne mois + annee numeriques - on construira le filtre avec MONTH et YEAR
        return None  # Cas special, gere separement
    elif grouping == "trimestre":
        return f"'T' + CAST(DATEPART(QUARTER, {f}) AS VARCHAR)"
    elif grouping == "trimestre_annee":
        return f"'T' + CAST(DATEPART(QUARTER, {f}) AS VARCHAR) + ' ' + CAST(YEAR({f}) AS VARCHAR)"
    elif grouping == "semestre":
        return f"CASE WHEN MONTH({f}) <= 6 THEN 'S1' ELSE 'S2' END"
    elif grouping == "semestre_annee":
        return f"CASE WHEN MONTH({f}) <= 6 THEN 'S1' ELSE 'S2' END + ' ' + CAST(YEAR({f}) AS VARCHAR)"
    elif grouping == "annee":
        return f"CAST(YEAR({f}) AS VARCHAR)"
    return None


def _build_date_grouping_where(field_name: str, grouping: str, value: str) -> Tuple[str, list]:
    """Construit la clause WHERE SQL pour un champ date avec regroupement"""
    f = f"[{field_name}]"

    if grouping == "jour":
        # value = "2022-08-15"
        return f"CONVERT(VARCHAR(10), {f}, 23) = ?", [value]

    elif grouping == "semaine":
        # value = "2022-S33"
        parts = value.split("-S")
        if len(parts) == 2:
            return f"YEAR({f}) = ? AND DATEPART(ISO_WEEK, {f}) = ?", [int(parts[0]), int(parts[1])]
        return f"1=1", []

    elif grouping == "mois":
        # value = "Aout"
        month_num = _MONTH_NAMES_FR_REVERSE.get(value)
        if month_num:
            return f"MONTH({f}) = ?", [month_num]
        return f"1=1", []

    elif grouping == "mois_annee":
        # value = "Aout 2022"
        parts = value.rsplit(" ", 1)
        if len(parts) == 2:
            month_num = _MONTH_NAMES_FR_REVERSE.get(parts[0])
            if month_num:
                return f"MONTH({f}) = ? AND YEAR({f}) = ?", [month_num, int(parts[1])]
        return f"1=1", []

    elif grouping == "trimestre":
        # value = "T3"
        q = value.replace("T", "")
        if q.isdigit():
            return f"DATEPART(QUARTER, {f}) = ?", [int(q)]
        return f"1=1", []

    elif grouping == "trimestre_annee":
        # value = "T3 2022"
        parts = value.split(" ")
        if len(parts) == 2:
            q = parts[0].replace("T", "")
            if q.isdigit() and parts[1].isdigit():
                return f"DATEPART(QUARTER, {f}) = ? AND YEAR({f}) = ?", [int(q), int(parts[1])]
        return f"1=1", []

    elif grouping == "semestre":
        # value = "S1"
        if value == "S1":
            return f"MONTH({f}) <= 6", []
        elif value == "S2":
            return f"MONTH({f}) > 6", []
        return f"1=1", []

    elif grouping == "semestre_annee":
        # value = "S1 2022"
        parts = value.split(" ")
        if len(parts) == 2 and parts[1].isdigit():
            if parts[0] == "S1":
                return f"MONTH({f}) <= 6 AND YEAR({f}) = ?", [int(parts[1])]
            elif parts[0] == "S2":
                return f"MONTH({f}) > 6 AND YEAR({f}) = ?", [int(parts[1])]
        return f"1=1", []

    elif grouping == "annee":
        # value = "2022"
        if value.isdigit():
            return f"YEAR({f}) = ?", [int(value)]
        return f"1=1", []

    # Pas de grouping - comparaison directe
    return f"{f} = ?", [value]


def _apply_date_grouping(value, grouping: str) -> str:
    """Transforme une valeur date selon le regroupement temporel demande"""
    if value is None:
        return ""
    # Parser la date si c'est une string
    dt = None
    if isinstance(value, datetime):
        dt = value
    elif isinstance(value, date):
        dt = datetime(value.year, value.month, value.day)
    elif isinstance(value, str):
        for fmt in ("%Y-%m-%d %H:%M:%S", "%Y-%m-%d", "%d/%m/%Y", "%Y-%m-%dT%H:%M:%S"):
            try:
                dt = datetime.strptime(value.split('.')[0], fmt)
                break
            except (ValueError, AttributeError):
                continue
    if dt is None:
        return str(value)

    if grouping == "jour":
        return dt.strftime("%Y-%m-%d")
    elif grouping == "semaine":
        iso = dt.isocalendar()
        return f"{iso[0]}-S{iso[1]:02d}"
    elif grouping == "mois":
        return _MONTH_NAMES_FR.get(dt.month, str(dt.month))
    elif grouping == "mois_annee":
        return f"{_MONTH_NAMES_FR.get(dt.month, str(dt.month))} {dt.year}"
    elif grouping == "trimestre":
        q = (dt.month - 1) // 3 + 1
        return f"T{q}"
    elif grouping == "trimestre_annee":
        q = (dt.month - 1) // 3 + 1
        return f"T{q} {dt.year}"
    elif grouping == "semestre":
        s = 1 if dt.month <= 6 else 2
        return f"S{s}"
    elif grouping == "semestre_annee":
        s = 1 if dt.month <= 6 else 2
        return f"S{s} {dt.year}"
    elif grouping == "annee":
        return str(dt.year)
    else:
        return str(value)


def _apply_date_groupings_to_data(
    source_data: List[Dict[str, Any]],
    row_fields: List[Dict[str, Any]],
    col_field_config: Optional[Dict[str, Any]] = None
) -> List[Dict[str, Any]]:
    """Applique tous les regroupements (date, numerique, texte) aux donnees source"""
    # Collecter tous les champs avec regroupement
    grouping_fields = {}  # field_name -> (type, config)

    all_fields = list(row_fields)
    if col_field_config:
        all_fields.append(col_field_config)

    for rf in all_fields:
        fname = rf.get("field", "")
        if rf.get("date_grouping"):
            grouping_fields[fname] = ("date", rf["date_grouping"])
        elif rf.get("numeric_grouping"):
            grouping_fields[fname] = ("numeric", rf["numeric_grouping"])
        elif rf.get("text_grouping"):
            grouping_fields[fname] = ("text", rf["text_grouping"])

    if not grouping_fields:
        return source_data

    result = []
    for row in source_data:
        new_row = dict(row)
        for field_name, (gtype, gconfig) in grouping_fields.items():
            if field_name in new_row:
                if gtype == "date":
                    new_row[field_name] = _apply_date_grouping(new_row[field_name], gconfig)
                elif gtype == "numeric":
                    new_row[field_name] = _apply_numeric_grouping(new_row[field_name], gconfig)
                elif gtype == "text":
                    new_row[field_name] = _apply_text_grouping(new_row[field_name], gconfig)
        result.append(new_row)
    return result


def _parse_date_sort_key(val: str):
    """Tente de parser une valeur pivot comme date pour le tri chronologique.
    Supporte: 'Janvier 2023', 'Janvier', 'T1 2023', 'S1 2023', '2023-S05', '2023-01-15', '2023', etc.
    Retourne un tuple (year, sub_order) pour le tri, ou None si non reconnu."""
    if not val:
        return None
    parts = val.strip().split()

    # Format: "Mois Annee" (ex: "Janvier 2023")
    if len(parts) == 2 and parts[0] in _MONTH_NAMES_FR_REVERSE:
        try:
            year = int(parts[1])
            month = _MONTH_NAMES_FR_REVERSE[parts[0]]
            return (year, month)
        except (ValueError, KeyError):
            pass

    # Format: "Mois" seul (ex: "Janvier")
    if len(parts) == 1 and parts[0] in _MONTH_NAMES_FR_REVERSE:
        return (0, _MONTH_NAMES_FR_REVERSE[parts[0]])

    # Format: "T1 2023" ou "T1"
    if parts[0].startswith('T') and len(parts[0]) == 2 and parts[0][1].isdigit():
        q = int(parts[0][1])
        year = int(parts[1]) if len(parts) == 2 and parts[1].isdigit() else 0
        return (year, q)

    # Format: "S1 2023" ou "S1"
    if parts[0].startswith('S') and len(parts[0]) == 2 and parts[0][1].isdigit():
        s = int(parts[0][1])
        year = int(parts[1]) if len(parts) == 2 and parts[1].isdigit() else 0
        return (year, s)

    # Format: "2023-S05" (semaine ISO)
    if '-S' in val:
        try:
            y, w = val.split('-S')
            return (int(y), int(w))
        except (ValueError, IndexError):
            pass

    # Format: "YYYY-MM-DD"
    if len(val) == 10 and val[4] == '-' and val[7] == '-':
        try:
            return (int(val[:4]), int(val[5:7]) * 100 + int(val[8:10]))
        except ValueError:
            pass

    # Format: annee seule "2023"
    if val.isdigit() and len(val) == 4:
        return (int(val), 0)

    return None


def _sort_pivot_columns(columns: List[str]) -> List[str]:
    """Trie les colonnes pivot intelligemment : chronologiquement si dates, sinon alphabetiquement."""
    if not columns:
        return columns

    # Tenter de parser toutes les colonnes comme dates
    parsed = []
    all_parsed = True
    for col in columns:
        key = _parse_date_sort_key(col)
        if key is not None:
            parsed.append((key, col))
        else:
            all_parsed = False
            break

    if all_parsed and parsed:
        # Tri chronologique
        parsed.sort(key=lambda x: x[0])
        return [col for _, col in parsed]

    # Fallback: tri alphabetique
    return sorted(columns)


def _build_pivot_data(
    source_data: List[Dict[str, Any]],
    row_fields: List[Dict[str, Any]],
    col_field: Optional[str],
    value_fields: List[Dict[str, Any]],
    show_subtotals: bool = True,
    show_grand_totals: bool = True,
    grand_total_position: str = "bottom",
    subtotal_position: str = "bottom",
    show_summary_row: bool = False,
    summary_functions: Optional[List[str]] = None
) -> Tuple[List[Dict[str, Any]], List[str]]:
    """Construit les donnees pivotees avec sous-totaux et totaux generaux"""
    if not source_data or not row_fields or not value_fields:
        return [], []

    row_field_names = [rf["field"] for rf in row_fields]

    # Extraire les valeurs uniques de la colonne pivot
    pivot_columns = []
    if col_field:
        pivot_cols_set = set()
        for row in source_data:
            val = row.get(col_field)
            if val is not None:
                pivot_cols_set.add(str(val))
        pivot_columns = _sort_pivot_columns(list(pivot_cols_set))

    # Construire les alias de valeurs
    value_aliases = []
    for vf in value_fields:
        agg = vf.get("aggregation", "SUM").upper()
        field = vf["field"]
        alias = f"{agg}_{field}"
        value_aliases.append({
            **vf,
            "alias": alias,
            "aggregation": agg
        })

    # Grouper les donnees
    groups = defaultdict(lambda: defaultdict(list))
    for row in source_data:
        row_key = tuple(str(row.get(f, "")) for f in row_field_names)
        col_val = str(row.get(col_field, "")) if col_field else "__NONE__"
        groups[row_key][col_val].append(row)

    # Calculer les agregations
    # Agregations non-additives (ne peuvent pas etre sommees pour sous-totaux)
    NON_ADDITIVE_AGGS = {"AVG", "VAR", "STDEV", "MEDIAN"}

    def aggregate(values: List[Any], agg: str) -> float:
        nums = [float(v) for v in values if v is not None and v != ""]
        if not nums:
            return 0
        if agg == "SUM":
            return sum(nums)
        elif agg == "COUNT":
            return len(nums)
        elif agg == "AVG":
            return sum(nums) / len(nums)
        elif agg == "MIN":
            return min(nums)
        elif agg == "MAX":
            return max(nums)
        elif agg == "DISTINCTCOUNT":
            return len(set(nums))
        elif agg == "VAR":
            mean = sum(nums) / len(nums)
            return sum((x - mean) ** 2 for x in nums) / len(nums)
        elif agg == "STDEV":
            mean = sum(nums) / len(nums)
            variance = sum((x - mean) ** 2 for x in nums) / len(nums)
            return variance ** 0.5
        elif agg == "MEDIAN":
            sorted_nums = sorted(nums)
            n = len(sorted_nums)
            mid = n // 2
            if n % 2 == 0:
                return (sorted_nums[mid - 1] + sorted_nums[mid]) / 2
            return sorted_nums[mid]
        return sum(nums)

    # Construire les lignes resultat
    result_rows = []
    grand_totals = defaultdict(float)
    grand_counts = defaultdict(int)

    # Grouper par premier champ de ligne pour les sous-totaux
    first_field_groups = defaultdict(list)

    for row_key in sorted(groups.keys()):
        # Ligne de donnees
        pivot_row = {
            "__rowKey__": "|".join(row_key),
            "__level__": len(row_field_names) - 1,
            "__isSubtotal__": False,
            "__isGrandTotal__": False,
        }

        # Champs de ligne
        for i, fname in enumerate(row_field_names):
            pivot_row[fname] = row_key[i]

        # Valeurs pivotees
        row_total_by_alias = defaultdict(float)

        if col_field and pivot_columns:
            for col_val in pivot_columns:
                cell_rows = groups[row_key].get(col_val, [])
                for va in value_aliases:
                    cell_values = [r.get(va["field"]) for r in cell_rows]
                    agg_val = aggregate(cell_values, va["aggregation"])
                    key = f"{col_val}__{va['alias']}"
                    pivot_row[key] = round(agg_val, 4)
                    row_total_by_alias[va["alias"]] += agg_val

                    # Grand totaux
                    grand_totals[f"{col_val}__{va['alias']}"] += agg_val
        else:
            # Pas de colonne pivot - agregation directe
            all_rows_for_key = []
            for col_val, rows in groups[row_key].items():
                all_rows_for_key.extend(rows)
            for va in value_aliases:
                cell_values = [r.get(va["field"]) for r in all_rows_for_key]
                agg_val = aggregate(cell_values, va["aggregation"])
                pivot_row[va["alias"]] = round(agg_val, 4)
                row_total_by_alias[va["alias"]] += agg_val
                grand_totals[va["alias"]] += agg_val

        # Totaux de ligne
        if col_field and pivot_columns:
            for va in value_aliases:
                pivot_row[f"__TOTAL__{va['alias']}"] = round(row_total_by_alias[va["alias"]], 4)
                grand_totals[f"__TOTAL__{va['alias']}"] += row_total_by_alias[va["alias"]]

        result_rows.append(pivot_row)

        # Tracking pour sous-totaux
        if len(row_field_names) > 1:
            first_val = row_key[0]
            first_field_groups[first_val].append(pivot_row)

    # Inserer les sous-totaux
    # Skip si chaque groupe n'a qu'une seule ligne (sous-total = doublon inutile)
    has_multi_row_groups = any(len(rows) > 1 for rows in first_field_groups.values()) if first_field_groups else False
    if show_subtotals and len(row_field_names) > 1 and has_multi_row_groups:
        rows_with_subtotals = []
        current_group = None

        def _insert_subtotal(grp_val):
            return _build_subtotal_row(
                first_field_groups[grp_val], row_field_names, grp_val,
                pivot_columns, value_aliases, col_field,
                raw_groups=groups, aggregate_fn=aggregate
            )

        for row in result_rows:
            first_val = row.get(row_field_names[0], "")

            if current_group is not None and first_val != current_group:
                # N'inserer le sous-total que si le groupe a plus d'une ligne
                if len(first_field_groups.get(current_group, [])) > 1:
                    st = _insert_subtotal(current_group)
                    rows_with_subtotals.append(st)

            if first_val != current_group and subtotal_position == "top":
                # Inserer sous-total AVANT le groupe (seulement si > 1 ligne)
                if first_val in first_field_groups and len(first_field_groups[first_val]) > 1:
                    st_top = _insert_subtotal(first_val)
                    rows_with_subtotals.append(st_top)

            current_group = first_val
            rows_with_subtotals.append(row)

        # Sous-total du dernier groupe
        if current_group is not None and current_group in first_field_groups and subtotal_position != "top":
            if len(first_field_groups[current_group]) > 1:
                rows_with_subtotals.append(_insert_subtotal(current_group))

        result_rows = rows_with_subtotals

    # Grand total
    if show_grand_totals:
        grand_row = {
            "__rowKey__": "__GRAND_TOTAL__",
            "__level__": 0,
            "__isSubtotal__": False,
            "__isGrandTotal__": True,
        }
        for fname in row_field_names:
            grand_row[fname] = "TOTAL"

        for va in value_aliases:
            if va.get("show_in_totals") is False:
                continue
            agg = (va.get("summary_aggregation") or va.get("aggregation", "SUM")).upper()
            if agg in NON_ADDITIVE_AGGS:
                # Re-agreger depuis les donnees brutes
                if col_field and pivot_columns:
                    for col_val in pivot_columns:
                        raw_vals = []
                        for rk, cols in groups.items():
                            raw_vals.extend([r.get(va["field"]) for r in cols.get(col_val, [])])
                        key = f"{col_val}__{va['alias']}"
                        grand_row[key] = round(aggregate(raw_vals, agg), 4)
                    # Total ligne
                    all_raw = []
                    for rk, cols in groups.items():
                        for cv, rows in cols.items():
                            all_raw.extend([r.get(va["field"]) for r in rows])
                    grand_row[f"__TOTAL__{va['alias']}"] = round(aggregate(all_raw, agg), 4)
                else:
                    all_raw = []
                    for rk, cols in groups.items():
                        for cv, rows in cols.items():
                            all_raw.extend([r.get(va["field"]) for r in rows])
                    grand_row[va["alias"]] = round(aggregate(all_raw, agg), 4)
            else:
                # Additive: utiliser les totaux accumules
                if col_field and pivot_columns:
                    for col_val in pivot_columns:
                        key = f"{col_val}__{va['alias']}"
                        grand_row[key] = round(grand_totals.get(key, 0), 4)
                    grand_row[f"__TOTAL__{va['alias']}"] = round(grand_totals.get(f"__TOTAL__{va['alias']}", 0), 4)
                else:
                    grand_row[va["alias"]] = round(grand_totals.get(va["alias"], 0), 4)

        if grand_total_position == "top":
            result_rows.insert(0, grand_row)
        else:
            result_rows.append(grand_row)

    # Ligne de resume (statistiques supplementaires)
    if show_summary_row and summary_functions:
        summary_row = {
            "__rowKey__": "__SUMMARY__",
            "__level__": 0,
            "__isSubtotal__": False,
            "__isGrandTotal__": False,
            "__isSummary__": True,
        }
        for fname in row_field_names:
            summary_row[fname] = "RESUME"

        for func in summary_functions:
            func_upper = func.upper()
            for va in value_aliases:
                if col_field and pivot_columns:
                    for col_val in pivot_columns:
                        raw_vals = []
                        for rk, cols in groups.items():
                            raw_vals.extend([r.get(va["field"]) for r in cols.get(col_val, [])])
                        key = f"{col_val}__{va['alias']}"
                        summary_row[f"__SUMM_{func_upper}__{key}"] = round(aggregate(raw_vals, func_upper), 4)
                    # Total
                    all_raw = []
                    for rk, cols in groups.items():
                        for cv, rows in cols.items():
                            all_raw.extend([r.get(va["field"]) for r in rows])
                    summary_row[f"__SUMM_{func_upper}____TOTAL__{va['alias']}"] = round(aggregate(all_raw, func_upper), 4)
                else:
                    all_raw = []
                    for rk, cols in groups.items():
                        for cv, rows in cols.items():
                            all_raw.extend([r.get(va["field"]) for r in rows])
                    summary_row[f"__SUMM_{func_upper}__{va['alias']}"] = round(aggregate(all_raw, func_upper), 4)

        result_rows.append(summary_row)

    return result_rows, pivot_columns


def _build_subtotal_row(
    group_rows: List[Dict[str, Any]],
    row_field_names: List[str],
    group_value: str,
    pivot_columns: List[str],
    value_aliases: List[Dict[str, Any]],
    col_field: Optional[str],
    raw_groups: Optional[Dict] = None,
    aggregate_fn=None
) -> Dict[str, Any]:
    """Construit une ligne de sous-total pour un groupe"""
    NON_ADDITIVE = {"AVG", "VAR", "STDEV", "MEDIAN"}

    subtotal = {
        "__rowKey__": f"__SUBTOTAL__{group_value}",
        "__level__": 0,
        "__isSubtotal__": True,
        "__isGrandTotal__": False,
        row_field_names[0]: f"Sous-total {group_value}",
    }
    for fname in row_field_names[1:]:
        subtotal[fname] = ""

    if col_field and pivot_columns:
        for col_val in pivot_columns:
            for va in value_aliases:
                if va.get("show_in_totals") is False:
                    continue
                key = f"{col_val}__{va['alias']}"
                agg = va.get("summary_aggregation") or va.get("aggregation", "SUM")
                if agg.upper() in NON_ADDITIVE and raw_groups and aggregate_fn:
                    # Re-agreger depuis les donnees brutes
                    raw_vals = []
                    for rk, cols in raw_groups.items():
                        if str(rk[0]) == str(group_value):
                            raw_vals.extend([r.get(va["field"]) for r in cols.get(col_val, [])])
                    subtotal[key] = round(aggregate_fn(raw_vals, agg.upper()), 4)
                else:
                    total = sum(r.get(key, 0) or 0 for r in group_rows)
                    subtotal[key] = round(total, 4)

        for va in value_aliases:
            if va.get("show_in_totals") is False:
                continue
            total_key = f"__TOTAL__{va['alias']}"
            agg = va.get("summary_aggregation") or va.get("aggregation", "SUM")
            if agg.upper() in NON_ADDITIVE and raw_groups and aggregate_fn:
                raw_vals = []
                for rk, cols in raw_groups.items():
                    if str(rk[0]) == str(group_value):
                        for cv, rows in cols.items():
                            raw_vals.extend([r.get(va["field"]) for r in rows])
                subtotal[total_key] = round(aggregate_fn(raw_vals, agg.upper()), 4)
            else:
                total = sum(r.get(total_key, 0) or 0 for r in group_rows)
                subtotal[total_key] = round(total, 4)
    else:
        for va in value_aliases:
            if va.get("show_in_totals") is False:
                continue
            agg = va.get("summary_aggregation") or va.get("aggregation", "SUM")
            if agg.upper() in NON_ADDITIVE and raw_groups and aggregate_fn:
                raw_vals = []
                for rk, cols in raw_groups.items():
                    if str(rk[0]) == str(group_value):
                        for cv, rows in cols.items():
                            raw_vals.extend([r.get(va["field"]) for r in rows])
                subtotal[va["alias"]] = round(aggregate_fn(raw_vals, agg.upper()), 4)
            else:
                total = sum(r.get(va["alias"], 0) or 0 for r in group_rows)
                subtotal[va["alias"]] = round(total, 4)

    return subtotal


def _add_percentages(
    pivot_data: List[Dict[str, Any]],
    pivot_columns: List[str],
    value_aliases: List[Dict[str, Any]],
    col_field: Optional[str],
    show_row_pct: bool,
    show_col_pct: bool,
    show_total_pct: bool
):
    """Ajoute les pourcentages (ligne, colonne, total) aux donnees pivotees en place"""
    if not pivot_data or not value_aliases:
        return

    # Trouver la ligne grand total pour les denominateurs colonne/total
    grand_row = None
    for row in pivot_data:
        if row.get("__isGrandTotal__"):
            grand_row = row
            break

    for row in pivot_data:
        if col_field and pivot_columns:
            for va in value_aliases:
                alias = va["alias"]
                row_total = row.get(f"__TOTAL__{alias}", 0) or 0
                grand_total = grand_row.get(f"__TOTAL__{alias}", 0) if grand_row else 0

                for col_val in pivot_columns:
                    key = f"{col_val}__{alias}"
                    val = row.get(key, 0) or 0

                    # % du total ligne
                    if show_row_pct:
                        row[f"{key}__pct_row"] = round(val / row_total * 100, 1) if row_total else 0

                    # % du total colonne
                    if show_col_pct:
                        col_total = grand_row.get(key, 0) if grand_row else 0
                        row[f"{key}__pct_col"] = round(val / col_total * 100, 1) if col_total else 0

                    # % du total general
                    if show_total_pct:
                        row[f"{key}__pct_total"] = round(val / grand_total * 100, 1) if grand_total else 0

                # Totaux de ligne aussi
                total_key = f"__TOTAL__{alias}"
                total_val = row.get(total_key, 0) or 0
                if show_total_pct:
                    row[f"{total_key}__pct_total"] = round(total_val / grand_total * 100, 1) if grand_total else 0
        else:
            # Pas de colonne pivot
            for va in value_aliases:
                alias = va["alias"]
                val = row.get(alias, 0) or 0
                grand_total = grand_row.get(alias, 0) if grand_row else 0

                if show_total_pct:
                    row[f"{alias}__pct_total"] = round(val / grand_total * 100, 1) if grand_total else 0


# =============================================================================
# NUMERIC / TEXT GROUPING
# =============================================================================

def _apply_numeric_grouping(value, config: Dict) -> str:
    """Regroupe une valeur numerique en intervalle."""
    if value is None:
        return ""
    try:
        num = float(value)
    except (ValueError, TypeError):
        return str(value)

    grouping_type = config.get("type", "interval")

    if grouping_type == "interval":
        step = config.get("step", 100)
        min_val = config.get("min", 0)
        bucket = int((num - min_val) // step) if step > 0 else 0
        low = min_val + bucket * step
        high = low + step
        return f"{int(low)}-{int(high)}"

    elif grouping_type == "ranges":
        for r in config.get("ranges", []):
            r_min = r.get("min", float('-inf'))
            r_max = r.get("max", float('inf'))
            if r_min <= num < r_max:
                return r.get("label", f"{r_min}-{r_max}")
        return config.get("other_label", "Autre")

    return str(value)


def _apply_text_grouping(value, config: Dict) -> str:
    """Regroupe une valeur texte par premiere lettre ou groupes custom."""
    if value is None:
        return ""
    s = str(value)
    grouping_type = config.get("type", "first_letter")

    if grouping_type == "first_letter":
        return s[0].upper() if s else ""

    elif grouping_type == "custom":
        groups = config.get("groups", {})
        for group_label, members in groups.items():
            if s in members:
                return group_label
        return config.get("other_label", "Autre")

    return s


# =============================================================================
# WINDOW CALCULATIONS
# =============================================================================

def _safe_eval_expression(expr_str: str, row: Dict, value_aliases: List[Dict]) -> float:
    """Evalue une expression de calcul de maniere securisee."""
    # Remplacer les references [NomChamp] par les valeurs
    field_refs = re.findall(r'\[([^\]]+)\]', expr_str)
    eval_expr = expr_str
    for ref in field_refs:
        ref_val = 0
        for va in value_aliases:
            if va.get("field") == ref or va.get("label") == ref or va.get("alias") == ref:
                ref_val = row.get(va["alias"], 0) or 0
                break
        eval_expr = eval_expr.replace(f"[{ref}]", str(float(ref_val)))

    # Securite : n'autoriser que les caracteres numeriques et operateurs
    cleaned = re.sub(r'[0-9\.\+\-\*/\(\)\s]', '', eval_expr)
    if cleaned:
        return 0  # Expression non securisee

    try:
        result = eval(eval_expr)
        return float(result)
    except Exception:
        return 0


def _apply_window_calculations(
    pivot_data: List[Dict[str, Any]],
    window_calcs: List[Dict[str, Any]],
    value_aliases: List[Dict[str, Any]],
    pivot_columns: List[str],
    col_field: Optional[str]
):
    """Applique les calculs de fenetre (cumul, difference, rang, expression)."""
    if not window_calcs:
        return

    # Filtrer les lignes de donnees (pas subtotals/grandtotal/summary)
    data_rows = [r for r in pivot_data
                 if not r.get("__isSubtotal__") and not r.get("__isGrandTotal__") and not r.get("__isSummary__")]

    for calc in window_calcs:
        calc_type = calc.get("type", "")
        source = calc.get("source_field", "")
        calc_id = calc.get("id") or f"_wc_{calc_type}_{source}"

        if calc_type == "running_total":
            running = 0
            for row in data_rows:
                val = row.get(source, 0) or 0
                running += float(val)
                row[f"__WC_{calc_id}__"] = round(running, 4)

        elif calc_type == "difference":
            prev = None
            for row in data_rows:
                val = float(row.get(source, 0) or 0)
                row[f"__WC_{calc_id}__"] = round(val - prev, 4) if prev is not None else 0
                prev = val

        elif calc_type == "pct_difference":
            prev = None
            for row in data_rows:
                val = float(row.get(source, 0) or 0)
                if prev is not None and prev != 0:
                    row[f"__WC_{calc_id}__"] = round((val - prev) / abs(prev) * 100, 2)
                else:
                    row[f"__WC_{calc_id}__"] = 0
                prev = val

        elif calc_type == "rank":
            vals = [(i, float(row.get(source, 0) or 0)) for i, row in enumerate(data_rows)]
            vals.sort(key=lambda x: x[1], reverse=True)
            for rank, (idx, _) in enumerate(vals, 1):
                data_rows[idx][f"__WC_{calc_id}__"] = rank

        elif calc_type == "expression":
            for row in data_rows:
                result = _safe_eval_expression(calc.get("expression", ""), row, value_aliases)
                row[f"__WC_{calc_id}__"] = round(result, calc.get("decimals", 4))


# =============================================================================
# CRUD ENDPOINTS
# =============================================================================

@router.get("")
async def list_pivots(
    user_id: Optional[int] = Query(None),
    is_public: Optional[bool] = Query(None)
):
    """Liste tous les pivots V2"""
    try:
        init_pivot_v2_tables()

        conditions = []
        params_list = []

        if user_id is not None:
            conditions.append("(created_by = ? OR is_public = 1)")
            params_list.append(user_id)
        if is_public is not None:
            conditions.append("is_public = ?")
            params_list.append(1 if is_public else 0)

        where = "WHERE " + " AND ".join(conditions) if conditions else ""
        query = f"""
            SELECT id, nom, description, data_source_id, data_source_code,
                   is_public, created_by, created_at, updated_at,
                   comparison_mode, application
            FROM APP_Pivots_V2
            {where}
            ORDER BY updated_at DESC
        """
        results = execute_query(query, tuple(params_list) if params_list else None, use_cache=False)
        return {"success": True, "data": results}
    except Exception as e:
        logger.error(f"Erreur liste pivots V2: {e}")
        return {"success": False, "error": str(e), "data": []}


@router.get("/{pivot_id}")
async def get_pivot(pivot_id: int):
    """Recupere la configuration complete d'un pivot"""
    try:
        init_pivot_v2_tables()

        query = """
            SELECT * FROM APP_Pivots_V2 WHERE id = ?
        """
        results = execute_query(query, (pivot_id,), use_cache=False)
        if not results:
            raise HTTPException(status_code=404, detail=f"Pivot {pivot_id} non trouve")

        row = results[0]
        # Parser les JSON
        row["rows_config"] = _parse_json(row.get("rows_config"), [])
        row["columns_config"] = _parse_json(row.get("columns_config"), [])
        row["filters_config"] = _parse_json(row.get("filters_config"), [])
        row["values_config"] = _parse_json(row.get("values_config"), [])
        row["formatting_rules"] = _parse_json(row.get("formatting_rules"), [])
        row["source_params"] = _parse_json(row.get("source_params"), [])

        return {"success": True, "data": row}
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Erreur get pivot V2 {pivot_id}: {e}")
        return {"success": False, "error": str(e)}


@router.post("")
async def create_pivot(data: PivotCreateRequest):
    """Cree un nouveau pivot"""
    try:
        init_pivot_v2_tables()

        # Auto-generer un code unique si absent
        code = getattr(data, 'code', None) or _generate_entity_code('PV', data.nom)

        with get_db_cursor() as cursor:
            cursor.execute("""
                INSERT INTO APP_Pivots_V2
                (nom, code, description, data_source_id, data_source_code,
                 rows_config, columns_config, filters_config, values_config,
                 show_grand_totals, show_subtotals, show_row_percent, show_col_percent,
                 show_total_percent, comparison_mode, formatting_rules, source_params,
                 grand_total_position, subtotal_position, show_summary_row, summary_functions,
                 window_calculations,
                 is_public, application, created_by, drilldown_data_source_code, created_at, updated_at)
                VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, GETDATE(), GETDATE())
            """, (
                data.nom,
                code,
                data.description,
                data.data_source_id,
                data.data_source_code,
                _to_json([f.dict() for f in data.rows_config] if data.rows_config else []),
                _to_json([f.dict() for f in data.columns_config] if data.columns_config else []),
                _to_json([f.dict() for f in data.filters_config] if data.filters_config else []),
                _to_json([f.dict() for f in data.values_config] if data.values_config else []),
                1 if data.show_grand_totals else 0,
                1 if data.show_subtotals else 0,
                1 if data.show_row_percent else 0,
                1 if data.show_col_percent else 0,
                1 if data.show_total_percent else 0,
                data.comparison_mode,
                _to_json([f.dict() for f in data.formatting_rules] if data.formatting_rules else []),
                _to_json(data.source_params or []),
                data.grand_total_position or "bottom",
                data.subtotal_position or "bottom",
                1 if data.show_summary_row else 0,
                _to_json(data.summary_functions),
                _to_json(data.window_calculations or []),
                1 if data.is_public else 0,
                data.application,
                data.created_by,
                data.drilldown_data_source_code,
            ))
            cursor.execute("SELECT @@IDENTITY AS id")
            new_id = cursor.fetchone()[0]

        return {"success": True, "message": "Pivot cree avec succes", "id": int(new_id)}
    except Exception as e:
        logger.error(f"Erreur creation pivot V2: {e}")
        return {"success": False, "error": str(e)}


@router.put("/{pivot_id}")
async def update_pivot(pivot_id: int, data: PivotUpdateRequest):
    """Met a jour un pivot existant"""
    try:
        # Construire la requete UPDATE dynamiquement
        updates = []
        params_list = []

        field_mapping = {
            "nom": data.nom,
            "description": data.description,
            "data_source_id": data.data_source_id,
            "data_source_code": data.data_source_code,
            "drilldown_data_source_code": data.drilldown_data_source_code,
            "show_grand_totals": (1 if data.show_grand_totals else 0) if data.show_grand_totals is not None else None,
            "show_subtotals": (1 if data.show_subtotals else 0) if data.show_subtotals is not None else None,
            "show_row_percent": (1 if data.show_row_percent else 0) if data.show_row_percent is not None else None,
            "show_col_percent": (1 if data.show_col_percent else 0) if data.show_col_percent is not None else None,
            "show_total_percent": (1 if data.show_total_percent else 0) if data.show_total_percent is not None else None,
            "grand_total_position": data.grand_total_position,
            "subtotal_position": data.subtotal_position,
            "show_summary_row": (1 if data.show_summary_row else 0) if data.show_summary_row is not None else None,
            "comparison_mode": data.comparison_mode,
            "is_public": (1 if data.is_public else 0) if data.is_public is not None else None,
            "application": data.application,
        }

        for field, value in field_mapping.items():
            if value is not None:
                updates.append(f"{field} = ?")
                params_list.append(value)

        # Champs JSON
        json_fields = {
            "rows_config": [f.dict() for f in data.rows_config] if data.rows_config is not None else None,
            "columns_config": [f.dict() for f in data.columns_config] if data.columns_config is not None else None,
            "filters_config": [f.dict() for f in data.filters_config] if data.filters_config is not None else None,
            "values_config": [f.dict() for f in data.values_config] if data.values_config is not None else None,
            "formatting_rules": [f.dict() for f in data.formatting_rules] if data.formatting_rules is not None else None,
            "source_params": data.source_params if data.source_params is not None else None,
            "summary_functions": data.summary_functions if data.summary_functions is not None else None,
            "window_calculations": data.window_calculations if data.window_calculations is not None else None,
        }

        for field, value in json_fields.items():
            if value is not None:
                updates.append(f"{field} = ?")
                params_list.append(_to_json(value))

        if not updates:
            return {"success": True, "message": "Rien a mettre a jour"}

        updates.append("updated_at = GETDATE()")
        params_list.append(pivot_id)

        query = f"UPDATE APP_Pivots_V2 SET {', '.join(updates)} WHERE id = ?"

        with get_db_cursor() as cursor:
            cursor.execute(query, tuple(params_list))

        return {"success": True, "message": "Pivot mis a jour avec succes"}
    except Exception as e:
        logger.error(f"Erreur update pivot V2 {pivot_id}: {e}")
        return {"success": False, "error": str(e)}


@router.delete("/{pivot_id}")
async def delete_pivot(pivot_id: int):
    """Supprime un pivot"""
    try:
        with get_db_cursor() as cursor:
            cursor.execute("DELETE FROM APP_Pivots_V2 WHERE id = ?", (pivot_id,))
        return {"success": True, "message": "Pivot supprime avec succes"}
    except Exception as e:
        logger.error(f"Erreur suppression pivot V2 {pivot_id}: {e}")
        return {"success": False, "error": str(e)}


# =============================================================================
# EXECUTION & PREVIEW
# =============================================================================

@router.post("/{pivot_id}/execute")
async def execute_pivot(
    pivot_id: int,
    request: PivotExecuteRequest,
    dwh_code:    Optional[str] = Header(None, alias="X-DWH-Code"),
    user_id_hdr: Optional[str] = Header(None, alias="X-User-Id"),
):
    """Execute le pivot et retourne les donnees pivotees"""
    from ..services.permissions import enforce_report_access
    try:
        uid = int(user_id_hdr) if user_id_hdr else None
        enforce_report_access(uid, 'pivot', pivot_id, dwh_code)
    except Exception as perm_err:
        if hasattr(perm_err, 'status_code'):
            raise

    start_time = time.time()

    try:
        # Charger la config du pivot
        results = execute_query("SELECT * FROM APP_Pivots_V2 WHERE id = ?", (pivot_id,), use_cache=False)
        if not results:
            raise HTTPException(status_code=404, detail=f"Pivot {pivot_id} non trouve")

        config = results[0]
        rows_config = _parse_json(config.get("rows_config"), [])
        columns_config = _parse_json(config.get("columns_config"), [])
        values_config = _parse_json(config.get("values_config"), [])
        formatting_rules = _parse_json(config.get("formatting_rules"), [])
        comparison_mode = config.get("comparison_mode")
        show_subtotals = bool(config.get("show_subtotals", True))
        show_grand_totals = bool(config.get("show_grand_totals", True))
        grand_total_position = config.get("grand_total_position", "bottom")
        subtotal_position = config.get("subtotal_position", "bottom")
        show_summary_row = bool(config.get("show_summary_row", False))
        summary_functions = _parse_json(config.get("summary_functions"), None)
        window_calculations = _parse_json(config.get("window_calculations"), [])

        context = request.context or {}

        # Mode raw : retourner les donnees brutes
        if request.raw:
            source_data, debug_info = _resolve_and_execute(config, context, dwh_code, limit=10000, return_debug=True)
            return {
                "success": True,
                "data": source_data,
                "rawMode": True,
                "total": len(source_data),
                "metadata": {
                    "executionTime": int((time.time() - start_time) * 1000)
                },
                "debug": debug_info
            }

        # Executer la requete source
        source_data, debug_info = _resolve_and_execute(config, context, dwh_code, limit=10000, return_debug=True)

        if not source_data:
            return {
                "success": True,
                "data": [],
                "pivotColumns": [],
                "rowFields": rows_config,
                "columnField": columns_config[0]["field"] if columns_config else None,
                "valueFields": values_config,
                "formattingRules": formatting_rules,
                "metadata": {"totalRows": 0, "executionTime": 0, "hasSubtotals": False},
                "debug": debug_info
            }

        # Champ colonne pivot (un seul supporte)
        col_field = columns_config[0]["field"] if columns_config else None

        # Appliquer les regroupements temporels
        col_field_config = columns_config[0] if columns_config else None
        source_data = _apply_date_groupings_to_data(source_data, rows_config, col_field_config)

        # Construire le pivot
        pivot_data, pivot_columns = _build_pivot_data(
            source_data,
            rows_config,
            col_field,
            values_config,
            show_subtotals=show_subtotals,
            show_grand_totals=show_grand_totals,
            grand_total_position=grand_total_position,
            subtotal_position=subtotal_position,
            show_summary_row=show_summary_row,
            summary_functions=summary_functions
        )

        # Construire les aliases avec labels
        value_fields_response = []
        for vf in values_config:
            agg = vf.get("aggregation", "SUM").upper()
            alias = f"{agg}_{vf['field']}"
            value_fields_response.append({
                "field": vf["field"],
                "aggregation": agg,
                "alias": alias,
                "label": vf.get("label", vf["field"]),
                "format": vf.get("format", "number"),
                "decimals": vf.get("decimals", 2),
            })

        # Calcul des pourcentages
        show_row_pct = bool(config.get("show_row_percent", False))
        show_col_pct = bool(config.get("show_col_percent", False))
        show_total_pct = bool(config.get("show_total_percent", False))
        if show_row_pct or show_col_pct or show_total_pct:
            _add_percentages(
                pivot_data, pivot_columns, value_fields_response,
                col_field, show_row_pct, show_col_pct, show_total_pct
            )

        # Calculs de fenetre (running total, difference, rank, expression)
        wc_response = []
        if window_calculations:
            _apply_window_calculations(pivot_data, window_calculations, value_fields_response, pivot_columns, col_field)
            wc_response = window_calculations

        # Comparaison N/N-1
        comparison_info = None
        if comparison_mode and pivot_columns and len(pivot_columns) >= 2:
            sorted_cols = sorted(pivot_columns, reverse=True)
            comparison_info = {
                "enabled": True,
                "mode": comparison_mode,
                "currentPeriod": sorted_cols[0],
                "previousPeriod": sorted_cols[1]
            }

        execution_time = int((time.time() - start_time) * 1000)

        return {
            "success": True,
            "data": pivot_data,
            "pivotColumns": pivot_columns,
            "rowFields": rows_config,
            "columnField": col_field,
            "valueFields": value_fields_response,
            "formattingRules": formatting_rules,
            "comparison": comparison_info,
            "windowCalculations": wc_response,
            "summaryFunctions": summary_functions if show_summary_row else None,
            "options": {
                "showGrandTotals": show_grand_totals,
                "showSubtotals": show_subtotals,
                "showRowPercent": bool(config.get("show_row_percent", False)),
                "showColPercent": bool(config.get("show_col_percent", False)),
                "showTotalPercent": bool(config.get("show_total_percent", False)),
                "showSummaryRow": show_summary_row,
            },
            "metadata": {
                "totalRows": len([r for r in pivot_data if not r.get("__isGrandTotal__")]),
                "sourceRows": len(source_data),
                "executionTime": execution_time,
                "hasSubtotals": show_subtotals and len(rows_config) > 1
            }
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Erreur execution pivot V2 {pivot_id}: {e}")
        return {"success": False, "error": str(e)}


@router.post("/{pivot_id}/preview")
async def preview_pivot(
    pivot_id: int,
    request: PivotExecuteRequest,
    dwh_code: Optional[str] = Header(None, alias="X-DWH-Code")
):
    """Apercu du pivot (limite a 100 lignes source)"""
    try:
        results = execute_query("SELECT * FROM APP_Pivots_V2 WHERE id = ?", (pivot_id,), use_cache=False)
        if not results:
            raise HTTPException(status_code=404, detail=f"Pivot {pivot_id} non trouve")

        config = results[0]
        context = request.context or {}

        # Executer avec limite
        source_data = _resolve_and_execute(config, context, dwh_code, limit=100)

        rows_config = _parse_json(config.get("rows_config"), [])
        columns_config = _parse_json(config.get("columns_config"), [])
        values_config = _parse_json(config.get("values_config"), [])

        col_field = columns_config[0]["field"] if columns_config else None

        # Appliquer les regroupements temporels
        col_field_config = columns_config[0] if columns_config else None
        source_data = _apply_date_groupings_to_data(source_data, rows_config, col_field_config)

        pivot_data, pivot_columns = _build_pivot_data(
            source_data, rows_config, col_field, values_config,
            show_subtotals=False, show_grand_totals=True
        )

        value_fields_response = []
        for vf in values_config:
            agg = vf.get("aggregation", "SUM").upper()
            value_fields_response.append({
                "field": vf["field"],
                "aggregation": agg,
                "alias": f"{agg}_{vf['field']}",
                "label": vf.get("label", vf["field"]),
                "format": vf.get("format", "number"),
                "decimals": vf.get("decimals", 2),
            })

        # Calcul des pourcentages pour preview
        show_row_pct = bool(config.get("show_row_percent", False))
        show_col_pct = bool(config.get("show_col_percent", False))
        show_total_pct = bool(config.get("show_total_percent", False))
        if show_row_pct or show_col_pct or show_total_pct:
            _add_percentages(
                pivot_data, pivot_columns, value_fields_response,
                col_field, show_row_pct, show_col_pct, show_total_pct
            )

        return {
            "success": True,
            "data": pivot_data,
            "pivotColumns": pivot_columns,
            "rowFields": rows_config,
            "columnField": col_field,
            "valueFields": value_fields_response,
            "options": {
                "showRowPercent": show_row_pct,
                "showColPercent": show_col_pct,
                "showTotalPercent": show_total_pct,
            },
            "preview": True,
            "sourceRows": len(source_data),
            "metadata": {"totalRows": len(pivot_data), "preview": True}
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Erreur preview pivot V2 {pivot_id}: {e}")
        return {"success": False, "error": str(e)}


# =============================================================================
# DRILL-DOWN
# =============================================================================

@router.post("/{pivot_id}/drilldown")
async def drilldown_pivot(
    pivot_id: int,
    request: PivotDrilldownRequest,
    dwh_code:    Optional[str] = Header(None, alias="X-DWH-Code"),
    user_id_hdr: Optional[str] = Header(None, alias="X-User-Id"),
):
    """Drill-down sur une cellule du pivot"""
    from ..services.permissions import enforce_report_access
    try:
        uid = int(user_id_hdr) if user_id_hdr else None
        enforce_report_access(uid, 'pivot', pivot_id, dwh_code)
    except Exception as perm_err:
        if hasattr(perm_err, 'status_code'):
            raise

    try:
        results = execute_query("SELECT * FROM APP_Pivots_V2 WHERE id = ?", (pivot_id,), use_cache=False)
        if not results:
            raise HTTPException(status_code=404, detail=f"Pivot {pivot_id} non trouve")

        config = results[0]
        rows_config = _parse_json(config.get("rows_config"), [])
        columns_config = _parse_json(config.get("columns_config"), [])
        context = request.context or {}

        # Construire un index des regroupements temporels
        date_groupings = {}
        for rf in rows_config:
            if rf.get("date_grouping"):
                date_groupings[rf["field"]] = rf["date_grouping"]
        if columns_config and columns_config[0].get("date_grouping"):
            date_groupings[columns_config[0]["field"]] = columns_config[0]["date_grouping"]

        # Recuperer la requete source (datasource drilldown dédiée si configurée, sinon principale)
        drill_ds_code = config.get("drilldown_data_source_code")
        ds_code = drill_ds_code or config.get("data_source_code")
        ds_id = config.get("data_source_id")

        if ds_code:
            ds = resolve_datasource(ds_code, dwh_code)
        elif ds_id:
            ds = resolve_datasource(ds_id, dwh_code)
        else:
            raise ValueError("Aucune source de donnees configuree")

        # Injecter les parametres de base
        base_query = inject_params(ds.query_template, context)
        base_query = re.sub(r'\s+ORDER\s+BY\s+[\s\S]+$', '', base_query, flags=re.IGNORECASE)

        # Construire les filtres drill-down
        where_parts = []
        filter_params = []

        for field_name, field_value in request.rowValues.items():
            if field_value is not None and field_value != "TOTAL" and not str(field_value).startswith("Sous-total"):
                if str(field_value) == "None" or field_value == "":
                    where_parts.append(f"[{field_name}] IS NULL")
                else:
                    grouping = date_groupings.get(field_name)
                    if grouping:
                        expr, params = _build_date_grouping_where(field_name, grouping, str(field_value))
                        where_parts.append(f"({expr})")
                        filter_params.extend(params)
                    else:
                        where_parts.append(f"[{field_name}] = ?")
                        filter_params.append(field_value)

        if request.columnValue:
            col_field = columns_config[0]["field"] if columns_config else None
            if col_field:
                grouping = date_groupings.get(col_field)
                if grouping:
                    expr, params = _build_date_grouping_where(col_field, grouping, str(request.columnValue))
                    where_parts.append(f"({expr})")
                    filter_params.extend(params)
                else:
                    where_parts.append(f"[{col_field}] = ?")
                    filter_params.append(request.columnValue)

        where_clause = " AND ".join(where_parts) if where_parts else "1=1"

        # Tri
        page_size = request.pageSize or 50
        page = max(1, request.page or 1)
        sort_field = request.sortField
        sort_direction = "DESC" if request.sortDirection and request.sortDirection.lower() == "desc" else "ASC"
        order_clause = f"ORDER BY [{sort_field}] {sort_direction}" if sort_field else "ORDER BY (SELECT NULL)"

        offset = (page - 1) * page_size

        # Requete COUNT
        count_query = f"SELECT COUNT(*) AS __cnt__ FROM ({base_query}) AS __src__ WHERE {where_clause}"

        # Requete paginee (SQL Server OFFSET/FETCH)
        drill_query = f"""
            SELECT * FROM ({base_query}) AS __src__
            WHERE {where_clause}
            {order_clause}
            OFFSET {offset} ROWS FETCH NEXT {page_size} ROWS ONLY
        """

        # Executer
        effective_dwh_dd = dwh_code
        if not effective_dwh_dd and ds.origin.value in ("template", "override"):
            try:
                dwh_list = execute_query("SELECT TOP 1 code FROM APP_DWH WHERE actif = 1 ORDER BY id", use_cache=True)
                if dwh_list:
                    effective_dwh_dd = dwh_list[0]['code']
            except Exception:
                pass

        params_tuple = tuple(filter_params) if filter_params else None

        if effective_dwh_dd:
            count_rows = DWHConnectionManager.execute_dwh_query(effective_dwh_dd, count_query, params_tuple, use_cache=False)
            page_data = DWHConnectionManager.execute_dwh_query(effective_dwh_dd, drill_query, params_tuple, use_cache=False)
        else:
            count_rows = execute_query(count_query, params_tuple, use_cache=False)
            page_data = execute_query(drill_query, params_tuple, use_cache=False)

        total = count_rows[0]["__cnt__"] if count_rows else 0
        total_pages = max(1, (total + page_size - 1) // page_size)
        page_data = [_normalize_row(r) for r in page_data]

        # Colonnes detectees
        columns = []
        numeric_fields = []
        if page_data:
            for col_name in page_data[0].keys():
                col_info = {"field": col_name, "header": col_name}
                sample = page_data[0].get(col_name)
                if isinstance(sample, (int, float)):
                    col_info["format"] = "number"
                    numeric_fields.append(col_name)
                elif isinstance(sample, (datetime, date)):
                    col_info["format"] = "date"
                columns.append(col_info)

        # Totaux des colonnes numeriques (sur toutes les lignes, pas juste la page)
        totals = {}
        if numeric_fields:
            sum_exprs = ", ".join(f"SUM([{f}]) AS [{f}]" for f in numeric_fields)
            totals_query = f"SELECT {sum_exprs} FROM ({base_query}) AS __src__ WHERE {where_clause}"
            if effective_dwh_dd:
                tot_rows = DWHConnectionManager.execute_dwh_query(effective_dwh_dd, totals_query, params_tuple, use_cache=False)
            else:
                tot_rows = execute_query(totals_query, params_tuple, use_cache=False)
            if tot_rows:
                row = _normalize_row(tot_rows[0])
                totals = {k: v for k, v in row.items() if v is not None}

        return {
            "success": True,
            "data": page_data,
            "total": total,
            "page": page,
            "pageSize": page_size,
            "totalPages": total_pages,
            "columns": columns,
            "totals": totals,
        }
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Erreur drill-down pivot V2 {pivot_id}: {e}")
        return {"success": False, "error": str(e)}


@router.post("/{pivot_id}/drilldown/export")
async def drilldown_pivot_export(
    pivot_id: int,
    request: PivotDrilldownRequest,
    dwh_code:    Optional[str] = Header(None, alias="X-DWH-Code"),
    user_id_hdr: Optional[str] = Header(None, alias="X-User-Id"),
):
    """Export Excel complet du drill-down (toutes les lignes)"""
    try:
        import pandas as pd

        results = execute_query("SELECT * FROM APP_Pivots_V2 WHERE id = ?", (pivot_id,), use_cache=False)
        if not results:
            raise HTTPException(status_code=404, detail=f"Pivot {pivot_id} non trouve")

        config = results[0]
        rows_config = _parse_json(config.get("rows_config"), [])
        columns_config = _parse_json(config.get("columns_config"), [])
        context = request.context or {}

        date_groupings = {}
        for rf in rows_config:
            if rf.get("date_grouping"):
                date_groupings[rf["field"]] = rf["date_grouping"]
        if columns_config and columns_config[0].get("date_grouping"):
            date_groupings[columns_config[0]["field"]] = columns_config[0]["date_grouping"]

        drill_ds_code = config.get("drilldown_data_source_code")
        ds_code = drill_ds_code or config.get("data_source_code")
        ds_id = config.get("data_source_id")

        if ds_code:
            ds = resolve_datasource(ds_code, dwh_code)
        elif ds_id:
            ds = resolve_datasource(ds_id, dwh_code)
        else:
            raise ValueError("Aucune source de donnees configuree")

        base_query = inject_params(ds.query_template, context)
        base_query = re.sub(r'\s+ORDER\s+BY\s+[\s\S]+$', '', base_query, flags=re.IGNORECASE)

        where_parts = []
        filter_params = []
        for field_name, field_value in request.rowValues.items():
            if field_value is not None and field_value != "TOTAL" and not str(field_value).startswith("Sous-total"):
                if str(field_value) == "None" or field_value == "":
                    where_parts.append(f"[{field_name}] IS NULL")
                else:
                    grouping = date_groupings.get(field_name)
                    if grouping:
                        expr, params = _build_date_grouping_where(field_name, grouping, str(field_value))
                        where_parts.append(f"({expr})")
                        filter_params.extend(params)
                    else:
                        where_parts.append(f"[{field_name}] = ?")
                        filter_params.append(field_value)

        if request.columnValue:
            col_field = columns_config[0]["field"] if columns_config else None
            if col_field:
                grouping = date_groupings.get(col_field)
                if grouping:
                    expr, params = _build_date_grouping_where(col_field, grouping, str(request.columnValue))
                    where_parts.append(f"({expr})")
                    filter_params.extend(params)
                else:
                    where_parts.append(f"[{col_field}] = ?")
                    filter_params.append(request.columnValue)

        where_clause = " AND ".join(where_parts) if where_parts else "1=1"

        sort_field = request.sortField
        sort_direction = "DESC" if request.sortDirection and request.sortDirection.lower() == "desc" else "ASC"
        order_clause = f"ORDER BY [{sort_field}] {sort_direction}" if sort_field else "ORDER BY (SELECT NULL)"

        export_query = f"""
            SELECT * FROM ({base_query}) AS __src__
            WHERE {where_clause}
            {order_clause}
        """

        effective_dwh_dd = dwh_code
        if not effective_dwh_dd and ds.origin.value in ("template", "override"):
            try:
                dwh_list = execute_query("SELECT TOP 1 code FROM APP_DWH WHERE actif = 1 ORDER BY id", use_cache=True)
                if dwh_list:
                    effective_dwh_dd = dwh_list[0]['code']
            except Exception:
                pass

        params_tuple = tuple(filter_params) if filter_params else None

        if effective_dwh_dd:
            all_data = DWHConnectionManager.execute_dwh_query(effective_dwh_dd, export_query, params_tuple, use_cache=False)
        else:
            all_data = execute_query(export_query, params_tuple, use_cache=False)

        all_data = [_normalize_row(r) for r in all_data]

        # Creer Excel en memoire
        df = pd.DataFrame(all_data) if all_data else pd.DataFrame()

        # Ligne de totaux
        if not df.empty:
            numeric_cols = df.select_dtypes(include='number').columns.tolist()
            totals_row = {col: df[col].sum() if col in numeric_cols else ('TOTAL' if col == df.columns[0] else '') for col in df.columns}
            df_total = pd.concat([df, pd.DataFrame([totals_row])], ignore_index=True)
        else:
            df_total = df

        buf = io.BytesIO()
        with pd.ExcelWriter(buf, engine='openpyxl') as writer:
            df_total.to_excel(writer, sheet_name='Detail', index=False)
            ws = writer.sheets['Detail']
            # Style en-tête
            from openpyxl.styles import Font, PatternFill, Alignment
            header_fill = PatternFill(start_color="1E3A5F", end_color="1E3A5F", fill_type="solid")
            total_fill  = PatternFill(start_color="EEF2FF", end_color="EEF2FF", fill_type="solid")
            for cell in ws[1]:
                cell.font = Font(bold=True, color="FFFFFF")
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal='center')
            # Style ligne totaux
            last_row = ws.max_row
            for cell in ws[last_row]:
                cell.font = Font(bold=True)
                cell.fill = total_fill
            # Largeur auto
            for col in ws.columns:
                max_len = max((len(str(c.value or '')) for c in col), default=10)
                ws.column_dimensions[col[0].column_letter].width = min(max_len + 2, 40)

        buf.seek(0)
        pivot_nom = config.get("nom", f"pivot_{pivot_id}").replace(" ", "_")
        filename = f"drilldown_{pivot_nom}_{datetime.now().strftime('%Y%m%d_%H%M%S')}.xlsx"

        return StreamingResponse(
            buf,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition": f"attachment; filename={filename}"}
        )
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"Erreur export drill-down pivot V2 {pivot_id}: {e}")
        raise HTTPException(status_code=500, detail=str(e))


# =============================================================================
# FIELDS
# =============================================================================

@router.get("/fields/{identifier}")
async def get_fields(
    identifier: str,
    dwh_code: Optional[str] = Header(None, alias="X-DWH-Code")
):
    """Recupere les champs disponibles d'une datasource"""
    try:
        # Identifier peut etre un ID numerique ou un code string
        if identifier.isdigit():
            ds = resolve_datasource(int(identifier), dwh_code)
        else:
            ds = resolve_datasource(identifier, dwh_code)

        # Injecter des valeurs par defaut pour obtenir les colonnes
        from ..services.parameter_resolver import get_default_context
        default_ctx = get_default_context()
        query = inject_params(ds.query_template, default_ctx)
        query = re.sub(r'\s+ORDER\s+BY\s+[\s\S]+$', '', query, flags=re.IGNORECASE)

        # Executer TOP 1 pour detecter les colonnes
        test_query = f"SELECT TOP 1 * FROM ({query}) AS __fields__"

        effective_dwh_f = dwh_code
        if not effective_dwh_f and ds.origin.value == "template":
            try:
                dwh_list = execute_query("SELECT TOP 1 code FROM APP_DWH WHERE actif = 1 ORDER BY id", use_cache=True)
                if dwh_list:
                    effective_dwh_f = dwh_list[0]['code']
            except Exception:
                pass

        if effective_dwh_f and ds.origin.value == "template":
            data = DWHConnectionManager.execute_dwh_query(effective_dwh_f, test_query, use_cache=False)
        else:
            data = execute_query(test_query, use_cache=False)

        fields = []
        if data:
            for key, value in data[0].items():
                field_type = "text"
                if isinstance(value, (int, float, Decimal)):
                    field_type = "number"
                elif isinstance(value, (datetime, date)):
                    field_type = "date"
                fields.append({"name": key, "type": field_type})

        return {"success": True, "fields": fields}
    except Exception as e:
        logger.error(f"Erreur get fields V2 '{identifier}': {e}")
        return {"success": False, "error": str(e), "fields": []}


# =============================================================================
# USER PREFERENCES
# =============================================================================

@router.get("/{pivot_id}/prefs/{user_id}")
async def get_user_prefs(pivot_id: int, user_id: int):
    """Recupere les preferences utilisateur d'un pivot"""
    try:
        init_pivot_v2_tables()
        query = """
            SELECT custom_config, ui_state, updated_at
            FROM APP_Pivot_User_Prefs
            WHERE pivot_id = ? AND user_id = ?
        """
        results = execute_query(query, (pivot_id, user_id), use_cache=False)

        if results:
            row = results[0]
            return {
                "success": True,
                "has_prefs": True,
                "data": {
                    "custom_config": _parse_json(row.get("custom_config"), {}),
                    "ui_state": _parse_json(row.get("ui_state"), {}),
                    "updated_at": row.get("updated_at")
                }
            }
        return {"success": True, "has_prefs": False, "data": None}
    except Exception as e:
        logger.error(f"Erreur get prefs V2: {e}")
        return {"success": False, "error": str(e)}


@router.put("/{pivot_id}/prefs/{user_id}")
async def save_user_prefs(pivot_id: int, user_id: int, request: PivotUserPrefsRequest):
    """Sauvegarde les preferences utilisateur (UPSERT)"""
    try:
        init_pivot_v2_tables()

        custom_config_json = _to_json(request.custom_config)
        ui_state_json = _to_json(request.ui_state)

        with get_db_cursor() as cursor:
            # Verifier si existe deja
            cursor.execute(
                "SELECT id FROM APP_Pivot_User_Prefs WHERE pivot_id = ? AND user_id = ?",
                (pivot_id, user_id)
            )
            existing = cursor.fetchone()

            if existing:
                cursor.execute("""
                    UPDATE APP_Pivot_User_Prefs
                    SET custom_config = ?, ui_state = ?, updated_at = GETDATE()
                    WHERE pivot_id = ? AND user_id = ?
                """, (custom_config_json, ui_state_json, pivot_id, user_id))
            else:
                cursor.execute("""
                    INSERT INTO APP_Pivot_User_Prefs (pivot_id, user_id, custom_config, ui_state, updated_at)
                    VALUES (?, ?, ?, ?, GETDATE())
                """, (pivot_id, user_id, custom_config_json, ui_state_json))

        return {"success": True, "message": "Preferences sauvegardees"}
    except Exception as e:
        logger.error(f"Erreur save prefs V2: {e}")
        return {"success": False, "error": str(e)}


@router.delete("/{pivot_id}/prefs/{user_id}")
async def reset_user_prefs(pivot_id: int, user_id: int):
    """Supprime les preferences utilisateur (reset)"""
    try:
        with get_db_cursor() as cursor:
            cursor.execute(
                "DELETE FROM APP_Pivot_User_Prefs WHERE pivot_id = ? AND user_id = ?",
                (pivot_id, user_id)
            )
        return {"success": True, "message": "Preferences reinitialises"}
    except Exception as e:
        logger.error(f"Erreur reset prefs V2: {e}")
        return {"success": False, "error": str(e)}


# =============================================================================
# EXPORT
# =============================================================================


def _build_export_headers(row_fields, pivot_columns, value_fields, col_field):
    """Construit la liste des en-tetes pour l'export."""
    headers = [rf.get("label", rf.get("field", "")) for rf in row_fields]
    if col_field and pivot_columns:
        for col_val in pivot_columns:
            for vf in value_fields:
                headers.append(f"{col_val} - {vf.get('label', vf['field'])}")
        for vf in value_fields:
            headers.append(f"Total - {vf.get('label', vf['field'])}")
    else:
        for vf in value_fields:
            headers.append(vf.get("label", vf["field"]))
    return headers


def _build_export_row(row, row_fields, pivot_columns, value_fields, col_field):
    """Extrait les valeurs d'une ligne pour l'export."""
    values = []
    for rf in row_fields:
        values.append(row.get(rf.get("field", ""), ""))

    if col_field and pivot_columns:
        for col_val in pivot_columns:
            for vf in value_fields:
                key = f"{col_val}__{vf['alias']}"
                values.append(row.get(key, 0))
        for vf in value_fields:
            key = f"__TOTAL__{vf['alias']}"
            values.append(row.get(key, 0))
    else:
        for vf in value_fields:
            values.append(row.get(vf["alias"], 0))

    return values


def _generate_excel(pivot_data, row_fields, pivot_columns, value_fields, col_field, title="Pivot"):
    """Genere un fichier Excel avec openpyxl."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    ws = wb.active
    ws.title = title[:31]  # Excel limite a 31 chars

    # Styles
    header_font = Font(name="Calibri", bold=True, size=11, color="FFFFFF")
    header_fill = PatternFill(start_color="2563EB", end_color="2563EB", fill_type="solid")
    title_font = Font(name="Calibri", bold=True, size=14, color="1E3A5F")
    subtotal_fill = PatternFill(start_color="E5E7EB", end_color="E5E7EB", fill_type="solid")
    subtotal_font = Font(name="Calibri", bold=True, size=11)
    grand_total_fill = PatternFill(start_color="1E40AF", end_color="1E40AF", fill_type="solid")
    grand_total_font = Font(name="Calibri", bold=True, size=11, color="FFFFFF")
    thin_border = Border(
        left=Side(style="thin", color="D1D5DB"),
        right=Side(style="thin", color="D1D5DB"),
        top=Side(style="thin", color="D1D5DB"),
        bottom=Side(style="thin", color="D1D5DB"),
    )
    number_format = '#,##0.00'

    # Titre
    ws.merge_cells(start_row=1, start_column=1, end_row=1, end_column=3)
    title_cell = ws.cell(row=1, column=1, value=title)
    title_cell.font = title_font
    ws.cell(row=2, column=1, value=f"Genere le {datetime.now().strftime('%d/%m/%Y %H:%M')}")
    ws.cell(row=2, column=1).font = Font(size=9, color="666666")

    # En-tetes (ligne 4)
    headers = _build_export_headers(row_fields, pivot_columns, value_fields, col_field)
    header_row = 4
    for col_idx, header in enumerate(headers, 1):
        cell = ws.cell(row=header_row, column=col_idx, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
        cell.border = thin_border

    # Donnees
    num_row_fields = len(row_fields)
    for data_idx, row in enumerate(pivot_data):
        excel_row = header_row + 1 + data_idx
        values = _build_export_row(row, row_fields, pivot_columns, value_fields, col_field)

        is_subtotal = row.get("__isSubtotal__", False)
        is_grand = row.get("__isGrandTotal__", False)

        for col_idx, val in enumerate(values, 1):
            cell = ws.cell(row=excel_row, column=col_idx, value=val)
            cell.border = thin_border

            # Format numerique pour les colonnes de valeurs
            if col_idx > num_row_fields and isinstance(val, (int, float, Decimal)):
                cell.number_format = number_format
                cell.alignment = Alignment(horizontal="right")
            else:
                cell.alignment = Alignment(horizontal="left")

            # Styles lignes speciales
            if is_grand:
                cell.fill = grand_total_fill
                cell.font = grand_total_font
            elif is_subtotal:
                cell.fill = subtotal_fill
                cell.font = subtotal_font

    # Largeur auto des colonnes
    for col_idx in range(1, len(headers) + 1):
        max_length = len(str(headers[col_idx - 1]))
        for data_row in range(header_row + 1, header_row + 1 + len(pivot_data)):
            cell_val = ws.cell(row=data_row, column=col_idx).value
            if cell_val:
                max_length = max(max_length, len(str(cell_val)))
        ws.column_dimensions[get_column_letter(col_idx)].width = min(max_length + 4, 40)

    # Freeze panes
    ws.freeze_panes = ws.cell(row=header_row + 1, column=num_row_fields + 1)

    output = io.BytesIO()
    wb.save(output)
    output.seek(0)
    return output


def _generate_pdf(pivot_data, row_fields, pivot_columns, value_fields, col_field, title="Pivot"):
    """Genere un fichier PDF avec reportlab."""
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4, landscape
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.units import cm

    output = io.BytesIO()
    doc = SimpleDocTemplate(output, pagesize=landscape(A4), topMargin=1 * cm, bottomMargin=1 * cm, leftMargin=1 * cm, rightMargin=1 * cm)
    elements = []
    styles = getSampleStyleSheet()

    # Titre
    title_style = ParagraphStyle("PivotTitle", parent=styles["Heading1"], fontSize=16, textColor=colors.HexColor("#1E3A5F"))
    elements.append(Paragraph(title, title_style))
    elements.append(Spacer(1, 0.3 * cm))
    date_style = ParagraphStyle("PivotDate", parent=styles["Normal"], fontSize=8, textColor=colors.grey)
    elements.append(Paragraph(f"Genere le {datetime.now().strftime('%d/%m/%Y %H:%M')}", date_style))
    elements.append(Spacer(1, 0.5 * cm))

    # En-tetes
    headers = _build_export_headers(row_fields, pivot_columns, value_fields, col_field)
    num_row_fields = len(row_fields)

    # Style cellules pour texte long
    cell_style = ParagraphStyle("Cell", parent=styles["Normal"], fontSize=7, leading=9)
    header_style = ParagraphStyle("Header", parent=styles["Normal"], fontSize=7, leading=9, textColor=colors.white, fontName="Helvetica-Bold")

    # Construire les donnees du tableau
    table_data = [[Paragraph(str(h), header_style) for h in headers]]

    for row in pivot_data:
        values = _build_export_row(row, row_fields, pivot_columns, value_fields, col_field)
        row_cells = []
        for col_idx, val in enumerate(values):
            if col_idx >= num_row_fields and isinstance(val, (int, float, Decimal)):
                row_cells.append(Paragraph(f"{val:,.2f}", cell_style))
            else:
                row_cells.append(Paragraph(str(val or ""), cell_style))
        table_data.append(row_cells)

    # Calculer les largeurs
    available = landscape(A4)[0] - 2 * cm
    num_cols = len(headers)
    col_widths = [available / num_cols] * num_cols

    table = Table(table_data, colWidths=col_widths, repeatRows=1)

    # Styles du tableau
    style_commands = [
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#2563EB")),
        ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
        ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE", (0, 0), (-1, 0), 7),
        ("FONTSIZE", (0, 1), (-1, -1), 7),
        ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#D1D5DB")),
        ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ("TOPPADDING", (0, 0), (-1, -1), 3),
        ("BOTTOMPADDING", (0, 0), (-1, -1), 3),
    ]

    # Alternance couleurs et styles sous-totaux/grand total
    for data_idx, row in enumerate(pivot_data):
        excel_row = data_idx + 1  # +1 pour header
        if row.get("__isGrandTotal__"):
            style_commands.append(("BACKGROUND", (0, excel_row), (-1, excel_row), colors.HexColor("#1E40AF")))
            style_commands.append(("TEXTCOLOR", (0, excel_row), (-1, excel_row), colors.white))
            style_commands.append(("FONTNAME", (0, excel_row), (-1, excel_row), "Helvetica-Bold"))
        elif row.get("__isSubtotal__"):
            style_commands.append(("BACKGROUND", (0, excel_row), (-1, excel_row), colors.HexColor("#E5E7EB")))
            style_commands.append(("FONTNAME", (0, excel_row), (-1, excel_row), "Helvetica-Bold"))
        elif data_idx % 2 == 0:
            style_commands.append(("BACKGROUND", (0, excel_row), (-1, excel_row), colors.HexColor("#F9FAFB")))

    table.setStyle(TableStyle(style_commands))
    elements.append(table)

    doc.build(elements)
    output.seek(0)
    return output


def _generate_pptx(pivot_data, row_fields, pivot_columns, value_fields, col_field, title="Pivot"):
    """Génère un fichier PowerPoint avec python-pptx."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    BLUE      = RGBColor(0x25, 0x63, 0xEB)
    DARK_BLUE = RGBColor(0x1E, 0x40, 0xAF)
    WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY_LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
    GRAY_ROW   = RGBColor(0xE5, 0xE7, 0xEB)

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── Slide 1 : titre ──────────────────────────────────────────────────────
    blank = prs.slide_layouts[6]  # layout totalement vide
    slide = prs.slides.add_slide(blank)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BLUE

    # Titre principal
    txb = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1.2))
    tf = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = title
    run.font.size = Pt(36)
    run.font.bold = True
    run.font.color.rgb = WHITE

    # Sous-titre date
    txb2 = slide.shapes.add_textbox(Inches(1), Inches(3.9), Inches(11.33), Inches(0.5))
    tf2 = txb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = f"Généré le {datetime.now().strftime('%d/%m/%Y à %H:%M')}  —  OptiBoard"
    r2.font.size = Pt(14)
    r2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)  # blue-200

    # ── Slides données (max 20 lignes / slide) ───────────────────────────────
    headers = _build_export_headers(row_fields, pivot_columns, value_fields, col_field)
    num_row_fields = len(row_fields)
    ROWS_PER_SLIDE = 20

    # Filtrer les lignes grand-total pour les garder en fin de chaque slide
    data_rows = [r for r in pivot_data if not r.get("__isGrandTotal__")]
    grand_rows = [r for r in pivot_data if r.get("__isGrandTotal__")]

    chunks = []
    for i in range(0, max(len(data_rows), 1), ROWS_PER_SLIDE):
        chunk = data_rows[i:i + ROWS_PER_SLIDE]
        if i + ROWS_PER_SLIDE >= len(data_rows):
            chunk = chunk + grand_rows
        chunks.append(chunk)

    if not chunks:
        chunks = [grand_rows or []]

    total_slides = len(chunks)
    for slide_idx, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(blank)

        # En-tête de slide
        hdr = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9), Inches(0.4))
        tf_h = hdr.text_frame
        p_h = tf_h.paragraphs[0]
        r_h = p_h.add_run()
        r_h.text = title
        r_h.font.size = Pt(11)
        r_h.font.bold = True
        r_h.font.color.rgb = DARK_BLUE

        if total_slides > 1:
            pag = slide.shapes.add_textbox(Inches(10), Inches(0.15), Inches(3), Inches(0.4))
            tf_p = pag.text_frame
            p_p = tf_p.paragraphs[0]
            p_p.alignment = PP_ALIGN.RIGHT
            r_p = p_p.add_run()
            r_p.text = f"{slide_idx + 1} / {total_slides}"
            r_p.font.size = Pt(9)
            r_p.font.color.rgb = RGBColor(0x6B, 0x72, 0x80)

        # Tableau
        num_cols = len(headers)
        num_rows_table = len(chunk) + 1  # +1 header

        if num_rows_table < 2:
            no_data = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
            no_data.text_frame.paragraphs[0].add_run().text = "Aucune donnée"
            continue

        left   = Inches(0.3)
        top    = Inches(0.65)
        width  = Inches(12.73)
        height = Inches(6.6)

        table = slide.shapes.add_table(num_rows_table, num_cols, left, top, width, height).table

        # Largeur colonnes : row_fields plus larges, value_fields plus étroites
        total_w = width
        rfw = Inches(1.8)
        vfw = (total_w - rfw * num_row_fields) // max(num_cols - num_row_fields, 1) if num_cols > num_row_fields else total_w // num_cols
        for ci in range(num_cols):
            table.columns[ci].width = rfw if ci < num_row_fields else vfw

        # Header row
        for ci, h in enumerate(headers):
            cell = table.cell(0, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = BLUE
            p_cell = cell.text_frame.paragraphs[0]
            p_cell.alignment = PP_ALIGN.CENTER
            run_c = p_cell.add_run()
            run_c.text = str(h)
            run_c.font.size = Pt(8)
            run_c.font.bold = True
            run_c.font.color.rgb = WHITE

        # Data rows
        for ri, row in enumerate(chunk):
            values = _build_export_row(row, row_fields, pivot_columns, value_fields, col_field)
            is_grand   = row.get("__isGrandTotal__", False)
            is_subtotal = row.get("__isSubtotal__", False)

            for ci, val in enumerate(values):
                cell = table.cell(ri + 1, ci)
                cell.fill.solid()
                if is_grand:
                    cell.fill.fore_color.rgb = DARK_BLUE
                elif is_subtotal:
                    cell.fill.fore_color.rgb = GRAY_ROW
                elif ri % 2 == 0:
                    cell.fill.fore_color.rgb = GRAY_LIGHT
                else:
                    cell.fill.fore_color.rgb = WHITE

                is_num = ci >= num_row_fields and isinstance(val, (int, float, Decimal))
                p_cell = cell.text_frame.paragraphs[0]
                p_cell.alignment = PP_ALIGN.RIGHT if is_num else PP_ALIGN.LEFT
                run_c = p_cell.add_run()
                run_c.text = f"{val:,.2f}" if is_num else str(val or "")
                run_c.font.size = Pt(8)
                if is_grand:
                    run_c.font.bold = True
                    run_c.font.color.rgb = WHITE
                elif is_subtotal:
                    run_c.font.bold = True

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output


@router.post("/{pivot_id}/export")
async def export_pivot(
    pivot_id: int,
    request: PivotExecuteRequest,
    format: str = Query("excel"),
    dwh_code:    Optional[str] = Header(None, alias="X-DWH-Code"),
    user_id_hdr: Optional[str] = Header(None, alias="X-User-Id"),
):
    """Export du pivot en Excel/CSV/PDF"""
    from ..services.permissions import enforce_report_access
    try:
        uid = int(user_id_hdr) if user_id_hdr else None
        enforce_report_access(uid, 'pivot', pivot_id, dwh_code)
    except Exception as perm_err:
        if hasattr(perm_err, 'status_code'):
            raise

    try:
        # Executer le pivot complet
        exec_result = await execute_pivot(pivot_id, request, dwh_code)

        if not exec_result.get("success"):
            return exec_result

        pivot_data = exec_result.get("data", [])
        value_fields = exec_result.get("valueFields", [])
        row_fields = exec_result.get("rowFields", [])
        pivot_columns = exec_result.get("pivotColumns", [])
        col_field = exec_result.get("columnField")

        # Recuperer le nom du pivot
        pivot_name = "Pivot"
        try:
            result = execute_query("SELECT nom FROM APP_Pivots_V2 WHERE id = ?", [pivot_id])
            if result:
                pivot_name = result[0].get("nom", "Pivot")
        except Exception:
            pass

        if format == "csv":
            headers = _build_export_headers(row_fields, pivot_columns, value_fields, col_field)
            csv_rows = [";".join(headers)]
            for row in pivot_data:
                values = _build_export_row(row, row_fields, pivot_columns, value_fields, col_field)
                csv_rows.append(";".join(str(v) for v in values))

            return {
                "success": True,
                "format": "csv",
                "content": "\n".join(csv_rows),
                "filename": f"pivot_{pivot_id}.csv"
            }

        if format == "excel":
            output = _generate_excel(pivot_data, row_fields, pivot_columns, value_fields, col_field, title=pivot_name)
            return StreamingResponse(
                output,
                media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                headers={"Content-Disposition": f'attachment; filename="pivot_{pivot_id}.xlsx"'}
            )

        if format == "pdf":
            output = _generate_pdf(pivot_data, row_fields, pivot_columns, value_fields, col_field, title=pivot_name)
            return StreamingResponse(
                output,
                media_type="application/pdf",
                headers={"Content-Disposition": f'attachment; filename="pivot_{pivot_id}.pdf"'}
            )

        if format == "pptx":
            output = _generate_pptx(pivot_data, row_fields, pivot_columns, value_fields, col_field, title=pivot_name)
            return StreamingResponse(
                output,
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": f'attachment; filename="pivot_{pivot_id}.pptx"'}
            )

        return {"success": False, "error": f"Format non supporte: {format}"}
    except Exception as e:
        logger.error(f"Erreur export pivot V2 {pivot_id}: {e}")
        return {"success": False, "error": str(e)}
