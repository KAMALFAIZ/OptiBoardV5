"""Routes pour le GridView Builder - Creation de vues grille personnalisees"""
import logging
from fastapi import APIRouter, HTTPException, Header
from pydantic import BaseModel
from typing import List, Optional, Dict, Any
from ..database_unified import (
    execute_central as execute_query,
    central_cursor as get_db_cursor,
    DWHConnectionManager,
    execute_central as execute_master_query,
    current_dwh_code as _ctx_dwh,
    execute_client,
    client_manager,
)
from ..services.datasource_resolver import datasource_resolver, DataSourceOrigin
import json
import re
import hashlib
import time

logger = logging.getLogger("GridViewBuilder")


# =============================================================================
# HELPER — Routage client DB / centrale pour APP_GridViews
# =============================================================================

def _gv_read(query: str, params: tuple = (), dwh_code: str = None) -> list:
    """
    Lit depuis APP_GridViews :
    - Si dwh_code fourni et base client dispo → base client (OptiBoard_cltXX)
    - Sinon → base centrale (OptiBoard_SaaS)
    """
    if dwh_code:
        try:
            if client_manager.has_client_db(dwh_code):
                return execute_client(query, params, dwh_code=dwh_code, use_cache=False)
        except Exception as e:
            logger.debug(f"_gv_read client fallback ({dwh_code}): {e}")
    return execute_query(query, params, use_cache=False)


def _gv_write(query: str, params: tuple = (), dwh_code: str = None) -> None:
    """
    Écrit dans APP_GridViews :
    - Si dwh_code fourni et base client dispo → base client
    - Sinon → base centrale
    """
    from ..database_unified import write_client, central_cursor
    if dwh_code:
        try:
            if client_manager.has_client_db(dwh_code):
                write_client(query, params, dwh_code=dwh_code)
                return
        except Exception as e:
            logger.debug(f"_gv_write client fallback ({dwh_code}): {e}")
    with get_db_cursor() as cursor:
        cursor.execute(query, params)


def _gv_cursor(dwh_code: str = None):
    """Retourne le bon cursor context manager selon le DWH code."""
    from ..database_unified import client_cursor
    if dwh_code:
        try:
            if client_manager.has_client_db(dwh_code):
                return client_cursor(dwh_code)
        except Exception:
            pass
    return get_db_cursor()


def _generate_entity_code(entity_type: str, nom: str) -> str:
    """Genere un code unique pour une entite (ex: GV_ventes_par_region_a3f2)"""
    slug = re.sub(r'[^a-z0-9]+', '_', nom.lower().strip())[:40].strip('_')
    suffix = hashlib.md5(f"{nom}{time.time()}".encode()).hexdigest()[:4]
    return f"{entity_type}_{slug}_{suffix}"

router = APIRouter(prefix="/api/gridview", tags=["gridview-builder"])


class ColumnConfig(BaseModel):
    field: str
    header: str
    width: Optional[int] = None
    sortable: bool = True
    filterable: bool = True
    format: Optional[str] = None  # number, currency, date, percent
    align: str = "left"  # left, center, right
    visible: bool = True
    pinned: Optional[str] = None  # 'left', 'right' ou null
    groupBy: bool = False  # Grouper par cette colonne par défaut

    class Config:
        extra = "allow"  # Permet _agState et autres champs frontend


class GridFeatures(BaseModel):
    """Options des fonctionnalités de la grille"""
    show_search: bool = True          # Recherche globale
    show_column_filters: bool = True  # Filtres par colonne
    show_grouping: bool = True        # Groupement
    show_column_toggle: bool = True   # Masquer/afficher colonnes
    show_export: bool = True          # Export CSV
    show_pagination: bool = True      # Pagination
    show_page_size: bool = True       # Sélecteur taille page
    allow_sorting: bool = True        # Tri des colonnes
    display_full_height: bool = True  # Colonnes 100% ajustees a la largeur du tableau


class GridViewConfig(BaseModel):
    id: Optional[int] = None
    nom: str
    description: Optional[str] = None
    data_source_id: Optional[int] = None  # Optionnel - peut être défini plus tard
    columns: List[ColumnConfig] = []
    default_sort: Optional[Dict[str, str]] = None  # {field: "CA", direction: "desc"}
    page_size: int = 25
    show_totals: bool = False
    total_columns: List[str] = []  # Colonnes a totaliser
    row_styles: List[Dict[str, Any]] = []  # Styles conditionnels
    is_public: bool = False
    features: Optional[GridFeatures] = None  # Options des fonctionnalités


class GridViewUpdate(BaseModel):
    nom: Optional[str] = None
    description: Optional[str] = None
    data_source_id: Optional[int] = None
    data_source_code: Optional[str] = None  # Nouveau: code du template DataSource
    columns: Optional[List[ColumnConfig]] = None
    default_sort: Optional[Dict[str, str]] = None
    default_group_by: Optional[List[str]] = None  # Champs de groupement par défaut
    page_size: Optional[int] = None
    show_totals: Optional[bool] = None
    total_columns: Optional[List[str]] = None
    row_styles: Optional[List[Dict[str, Any]]] = None
    is_public: Optional[bool] = None
    application: Optional[str] = None
    features: Optional[GridFeatures] = None


def init_gridview_tables():
    """Cree les tables pour le gridview builder"""
    # Créer la table si elle n'existe pas
    create_query = """
    IF NOT EXISTS (SELECT * FROM sysobjects WHERE name='APP_GridViews' AND xtype='U')
    CREATE TABLE APP_GridViews (
        id INT IDENTITY(1,1) PRIMARY KEY,
        nom NVARCHAR(200) NOT NULL,
        description NVARCHAR(500),
        data_source_id INT,
        columns_config NVARCHAR(MAX),
        default_sort NVARCHAR(200),
        page_size INT DEFAULT 25,
        show_totals BIT DEFAULT 0,
        total_columns NVARCHAR(MAX),
        row_styles NVARCHAR(MAX),
        features NVARCHAR(MAX),
        is_public BIT DEFAULT 0,
        created_by INT,
        created_at DATETIME DEFAULT GETDATE(),
        updated_at DATETIME DEFAULT GETDATE()
    )
    """

    # Liste des colonnes à ajouter si elles n'existent pas
    columns_to_add = [
        ("code", "VARCHAR(100)"),  # Code unique pour publication Master
        ("data_source_id", "INT"),
        ("data_source_code", "VARCHAR(100)"),  # Nouveau: code du template DataSource
        ("columns_config", "NVARCHAR(MAX)"),
        ("default_sort", "NVARCHAR(200)"),
        ("page_size", "INT DEFAULT 25"),
        ("show_totals", "BIT DEFAULT 0"),
        ("total_columns", "NVARCHAR(MAX)"),
        ("row_styles", "NVARCHAR(MAX)"),
        ("features", "NVARCHAR(MAX)"),
        ("default_group_by", "NVARCHAR(MAX)"),  # Champs de groupement par défaut (JSON array)
        ("application", "NVARCHAR(100)"),
        ("is_public", "BIT DEFAULT 0"),
        ("created_by", "INT"),
        ("created_at", "DATETIME DEFAULT GETDATE()"),
        ("updated_at", "DATETIME DEFAULT GETDATE()"),
    ]

    try:
        with get_db_cursor() as cursor:
            cursor.execute(create_query)

            # Ajouter les colonnes manquantes
            for col_name, col_type in columns_to_add:
                add_col_query = f"""
                IF NOT EXISTS (SELECT * FROM sys.columns WHERE object_id = OBJECT_ID('APP_GridViews') AND name = '{col_name}')
                ALTER TABLE APP_GridViews ADD {col_name} {col_type}
                """
                cursor.execute(add_col_query)
        return True
    except Exception as e:
        print(f"Erreur init gridview tables: {e}")
        return False


def serialize_model(obj):
    """Sérialise un objet Pydantic (compatible v1 et v2)"""
    if obj is None:
        return None
    if hasattr(obj, 'model_dump'):
        return obj.model_dump()
    elif hasattr(obj, 'dict'):
        return obj.dict()
    return obj


@router.get("/grids")
async def get_gridviews(
    user_id: Optional[int] = None,
    dwh_code: Optional[str] = Header(None, alias="X-DWH-Code"),
):
    """Liste toutes les grilles (depuis base client si X-DWH-Code fourni, sinon centrale)"""
    try:
        init_gridview_tables()

        if user_id:
            results = _gv_read(
                """SELECT id, nom, code, description, data_source_id, page_size, is_public, created_by, created_at, updated_at, application
                   FROM APP_GridViews
                   WHERE created_by = ? OR is_public = 1
                   ORDER BY updated_at DESC""",
                (user_id,),
                dwh_code=dwh_code,
            )
        else:
            results = _gv_read(
                """SELECT id, nom, code, description, data_source_id, page_size, is_public, created_by, created_at, updated_at, application
                   FROM APP_GridViews
                   ORDER BY updated_at DESC""",
                dwh_code=dwh_code,
            )
        return {"success": True, "data": results}
    except Exception as e:
        return {"success": False, "error": str(e), "data": []}


@router.get("/grids/{grid_id}")
async def get_gridview(
    grid_id: int,
    dwh_code: Optional[str] = Header(None, alias="X-DWH-Code"),
):
    """Recupere une grille avec sa configuration (depuis base client si X-DWH-Code fourni)"""
    try:
        # S'assurer que la colonne features existe
        init_gridview_tables()

        results = _gv_read(
            "SELECT * FROM APP_GridViews WHERE id = ?",
            (grid_id,),
            dwh_code=dwh_code,
        )
        if not results:
            raise HTTPException(status_code=404, detail="Grille non trouvee")

        grid = results[0]
        # Parser les configs JSON
        grid['columns'] = json.loads(grid.get('columns_config') or '[]')
        grid['default_sort'] = json.loads(grid.get('default_sort') or 'null')
        grid['total_columns'] = json.loads(grid.get('total_columns') or '[]')
        grid['default_group_by'] = json.loads(grid.get('default_group_by') or '[]')
        grid['row_styles'] = json.loads(grid.get('row_styles') or '[]')
        # Features avec valeurs par défaut
        default_features = {
            'show_search': True,
            'show_column_filters': True,
            'show_grouping': True,
            'show_column_toggle': True,
            'show_export': True,
            'show_pagination': True,
            'show_page_size': True,
            'allow_sorting': True,
            'display_full_height': True
        }
        raw_features = grid.get('features')
        saved_features = json.loads(raw_features or '{}')
        grid['features'] = {**default_features, **saved_features}

        return {"success": True, "data": grid}
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/grids")
async def create_gridview(grid: Dict[str, Any], user_id: int = 1):
    """Cree une nouvelle grille"""
    try:
        init_gridview_tables()

        nom = grid.get('nom', 'Sans nom')
        description = grid.get('description', '')
        data_source_id = grid.get('data_source_id')
        data_source_code = grid.get('data_source_code')  # Nouveau: code du template
        columns = grid.get('columns', [])
        default_sort = grid.get('default_sort')
        page_size = grid.get('page_size', 25)
        show_totals = grid.get('show_totals', False)
        total_columns = grid.get('total_columns', [])
        row_styles = grid.get('row_styles', [])
        features = grid.get('features', {})
        is_public = grid.get('is_public', False)

        # Auto-generer un code unique si absent
        code = grid.get('code') or _generate_entity_code('GV', nom)

        # Convertir data_source_id: 0 en None
        if data_source_id is not None and data_source_id <= 0:
            data_source_id = None

        with get_db_cursor() as cursor:
            cursor.execute(
                """INSERT INTO APP_GridViews (nom, code, description, data_source_id, data_source_code, columns_config, default_sort, page_size, show_totals, total_columns, row_styles, features, is_public, created_by)
                   VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)""",
                (
                    nom,
                    code,
                    description,
                    data_source_id,
                    data_source_code,
                    json.dumps(columns),
                    json.dumps(default_sort),
                    page_size,
                    show_totals,
                    json.dumps(total_columns),
                    json.dumps(row_styles),
                    json.dumps(features),
                    is_public,
                    user_id
                )
            )
            cursor.execute("SELECT @@IDENTITY AS id")
            new_id = cursor.fetchone()[0]

        return {"success": True, "message": "Grille creee", "id": new_id}
    except Exception as e:
        print(f"Erreur creation gridview: {e}")
        raise HTTPException(status_code=400, detail=str(e))


@router.put("/grids/{grid_id}")
async def update_gridview(grid_id: int, grid: GridViewUpdate):
    """Met a jour une grille"""
    try:
        # S'assurer que la colonne features existe
        init_gridview_tables()

        updates = []
        params = []

        if grid.nom is not None:
            updates.append("nom = ?")
            params.append(grid.nom)
        if grid.description is not None:
            updates.append("description = ?")
            params.append(grid.description)
        if grid.data_source_id is not None:
            updates.append("data_source_id = ?")
            params.append(grid.data_source_id)
        if grid.data_source_code is not None:
            updates.append("data_source_code = ?")
            params.append(grid.data_source_code)
        if grid.columns is not None:
            updates.append("columns_config = ?")
            params.append(json.dumps([serialize_model(c) for c in grid.columns]))
        if grid.default_sort is not None:
            updates.append("default_sort = ?")
            params.append(json.dumps(grid.default_sort))
        if grid.page_size is not None:
            updates.append("page_size = ?")
            params.append(grid.page_size)
        if grid.show_totals is not None:
            updates.append("show_totals = ?")
            params.append(grid.show_totals)
        if grid.total_columns is not None:
            updates.append("total_columns = ?")
            params.append(json.dumps(grid.total_columns))
        if grid.default_group_by is not None:
            updates.append("default_group_by = ?")
            params.append(json.dumps(grid.default_group_by))
        if grid.row_styles is not None:
            updates.append("row_styles = ?")
            params.append(json.dumps([serialize_model(s) if hasattr(s, 'dict') or hasattr(s, 'model_dump') else s for s in grid.row_styles]))
        if grid.is_public is not None:
            updates.append("is_public = ?")
            params.append(grid.is_public)
        if grid.application is not None:
            updates.append("application = ?")
            params.append(grid.application)

        # Sauvegarder features si présent
        if grid.features is not None:
            features_data = serialize_model(grid.features)
            updates.append("features = ?")
            params.append(json.dumps(features_data))

        updates.append("updated_at = GETDATE()")

        if not updates:
            return {"success": False, "message": "Aucune modification"}

        params.append(grid_id)

        query = f"UPDATE APP_GridViews SET {', '.join(updates)} WHERE id = ?"

        with get_db_cursor() as cursor:
            cursor.execute(query, tuple(params))

        return {"success": True, "message": "Grille mise a jour"}
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))


@router.delete("/grids/{grid_id}")
async def delete_gridview(grid_id: int):
    """Supprime une grille"""
    try:
        with get_db_cursor() as cursor:
            cursor.execute("DELETE FROM APP_GridViews WHERE id = ?", (grid_id,))
        return {"success": True, "message": "Grille supprimee"}
    except Exception as e:
        raise HTTPException(status_code=400, detail=str(e))


class GridDataRequest(BaseModel):
    """Requête pour récupérer les données d'une grille"""
    page: int = 1
    page_size: Optional[int] = None
    sort_field: Optional[str] = None
    sort_direction: str = "asc"
    filters: Dict[str, Any] = {}
    context: Dict[str, Any] = {}  # Contexte des filtres globaux


@router.post("/grids/{grid_id}/data")
async def get_grid_data(
    grid_id: int,
    request: GridDataRequest = GridDataRequest(),
    dwh_code:    Optional[str] = Header(None, alias="X-DWH-Code"),
    user_id_hdr: Optional[str] = Header(None, alias="X-User-Id"),
):
    """Recupere les donnees d'une grille avec pagination, tri et filtres globaux

    Utilise le resolver avec fallback Template -> Override pour les DataSources

    Args:
        grid_id: ID de la grille
        request: Configuration incluant le contexte des filtres globaux
        dwh_code: Code du DWH pour la resolution des DataSources
    """
    from ..services.parameter_resolver import inject_params
    from ..services.permissions import enforce_report_access

    # ── Vérification de permission ──────────────────────────────────────────
    try:
        uid = int(user_id_hdr) if user_id_hdr else None
        enforce_report_access(uid, 'gridview', grid_id, dwh_code)
    except Exception as perm_err:
        from fastapi import HTTPException as _HTTPException
        if hasattr(perm_err, 'status_code'):
            raise
        # Erreur technique dans le checker → fail-open, on continue

    try:
        # Recuperer la config de la grille (base client si dwh_code, sinon centrale)
        results = _gv_read(
            "SELECT * FROM APP_GridViews WHERE id = ?",
            (grid_id,),
            dwh_code=dwh_code,
        )
        if not results:
            raise HTTPException(status_code=404, detail="Grille non trouvee")

        grid = results[0]
        columns = json.loads(grid.get('columns_config') or '[]')
        default_sort = json.loads(grid.get('default_sort') or 'null')
        grid_page_size = request.page_size or grid.get('page_size', 25)
        show_totals = grid.get('show_totals', False)
        total_columns = json.loads(grid.get('total_columns') or '[]')

        # Recuperer la source de donnees avec fallback Template -> Override
        datasource_id = grid.get('data_source_id')
        datasource_code = grid.get('data_source_code')
        base_query = None
        origin = None

        try:
            if datasource_code:
                # Resolution par code (prioritaire)
                datasource = datasource_resolver.resolve_by_code(datasource_code, dwh_code)
                base_query = datasource.query_template
                origin = datasource.origin.value
            elif datasource_id:
                # Resolution par ID
                datasource = datasource_resolver.resolve_by_id(datasource_id, dwh_code)
                base_query = datasource.query_template
                origin = datasource.origin.value
            else:
                raise HTTPException(status_code=404, detail="Source de donnees non configuree")
        except ValueError:
            # Fallback: ancien comportement
            if datasource_id:
                source = execute_query(
                    "SELECT * FROM APP_DataSources WHERE id = ?",
                    (datasource_id,),
                    use_cache=False
                )
                if not source:
                    raise HTTPException(status_code=404, detail="Source de donnees non trouvee")
                base_query = source[0]['query_template']
            else:
                raise HTTPException(status_code=404, detail="Source de donnees non trouvee")

        # Gerer @societe_filter AVANT inject_params
        if '@societe_filter' in base_query:
            societe_val = request.context.get('societe') or request.context.get('societe_filter')
            if societe_val:
                base_query = base_query.replace('@societe_filter', f"societe = '{societe_val}'")
            else:
                base_query = base_query.replace('@societe_filter', '1=1')

        # Injecter les paramètres directement depuis le contexte
        final_query = inject_params(base_query, request.context)

        # Router vers le DWH pour les datasources template
        effective_dwh = dwh_code or _ctx_dwh.get()
        if not effective_dwh:
            try:
                dwh_list = execute_master_query(
                    "SELECT TOP 1 code FROM APP_DWH WHERE actif = 1 ORDER BY id",
                    use_cache=True
                )
                if dwh_list:
                    effective_dwh = dwh_list[0]['code']
            except Exception:
                pass

        if effective_dwh:
            try:
                all_data = DWHConnectionManager.execute_dwh_query(
                    effective_dwh, final_query, use_cache=False
                )
                logger.debug(f"[GRID {grid_id}] DWH={effective_dwh} → {len(all_data)} lignes")
            except Exception as dwh_err:
                logger.error(f"[GRID {grid_id}] Erreur connexion DWH '{effective_dwh}': {dwh_err}")
                return {
                    "success": False,
                    "error": f"Connexion DWH '{effective_dwh}' impossible : {type(dwh_err).__name__} — {str(dwh_err)[:200]}",
                    "data": [], "total": 0, "page": 1,
                    "page_size": grid_page_size, "total_pages": 0,
                    "totals": {}, "columns": columns
                }
        else:
            all_data = execute_query(final_query, use_cache=False)

        if not all_data:
            logger.info(f"[GRID {grid_id}] 0 lignes — DWH={effective_dwh} — Query: {final_query[:300]}")
            return {
                "success": True,
                "data": [], "total": 0, "page": 1,
                "page_size": grid_page_size, "total_pages": 0,
                "totals": {}, "columns": columns,
            }

        # ── Pagination SQL si toutes les lignes ont été chargées ─────────────
        # Optimisation : si plus de 500 lignes, utiliser pagination SQL directe
        total_count = len(all_data)

        # Tri
        sort = request.sort_field or (default_sort.get('field') if default_sort else None)
        direction = request.sort_direction if request.sort_field else (
            default_sort.get('direction', 'asc') if default_sort else 'asc'
        )

        if total_count > 500 and effective_dwh:
            # ── Pagination SQL native (OFFSET / FETCH NEXT) ──────────────────
            # On relance 2 requêtes : COUNT(*) + données paginées
            try:
                # Wrapper COUNT
                count_query = f"SELECT COUNT(*) AS total FROM ({final_query}) AS _cnt_sub"
                cnt_res = DWHConnectionManager.execute_dwh_query(effective_dwh, count_query, use_cache=False)
                total_count = cnt_res[0]["total"] if cnt_res else total_count

                # Trier + paginer en SQL
                safe_sort = re.sub(r'[^a-zA-Z0-9_ \[\]éèàùâêîôûäëïöüç]', '', sort) if sort else None
                order_clause = f'ORDER BY [{safe_sort}] {"DESC" if direction.lower()=="desc" else "ASC"}' if safe_sort else "ORDER BY (SELECT NULL)"
                offset_val = (request.page - 1) * grid_page_size
                paginated_query = (
                    f"SELECT * FROM ({final_query}) AS _paged_sub\n"
                    f"{order_clause}\n"
                    f"OFFSET {offset_val} ROWS FETCH NEXT {grid_page_size} ROWS ONLY"
                )
                data = DWHConnectionManager.execute_dwh_query(effective_dwh, paginated_query, use_cache=False)

                # Totaux (seulement si demandé — requête COUNT séparée par colonne)
                totals = {}
                if show_totals and total_columns:
                    agg = ", ".join(f"SUM(CAST([{c}] AS FLOAT)) AS [{c}]" for c in total_columns)
                    tot_q = f"SELECT {agg} FROM ({final_query}) AS _tot_sub"
                    try:
                        tot_res = DWHConnectionManager.execute_dwh_query(effective_dwh, tot_q, use_cache=False)
                        totals = tot_res[0] if tot_res else {}
                    except Exception:
                        totals = {}

                logger.debug(f"[GRID {grid_id}] SQL-paginated: total={total_count}, page={request.page}, rows={len(data)}")

            except Exception as pg_err:
                # Fallback pagination Python si SQL pagination échoue
                logger.warning(f"[GRID {grid_id}] SQL pagination failed, fallback Python: {pg_err}")
                if sort and all_data and sort in all_data[0]:
                    reverse = direction.lower() == 'desc'
                    all_data = sorted(all_data, key=lambda x: (x.get(sort) is None, x.get(sort) or 0), reverse=reverse)
                offset = (request.page - 1) * grid_page_size
                data = all_data[offset:offset + grid_page_size]
                totals = {}
                if show_totals and total_columns:
                    for col in total_columns:
                        try:
                            totals[col] = sum(r.get(col, 0) or 0 for r in all_data if isinstance(r.get(col), (int, float)))
                        except Exception:
                            totals[col] = 0
        else:
            # ── Pagination Python (≤500 lignes ou pas de DWH) ────────────────
            if sort and all_data and sort in all_data[0]:
                reverse = direction.lower() == 'desc'
                all_data = sorted(all_data, key=lambda x: (x.get(sort) is None, x.get(sort) or 0), reverse=reverse)
            offset = (request.page - 1) * grid_page_size
            data = all_data[offset:offset + grid_page_size]
            totals = {}
            if show_totals and total_columns:
                for col in total_columns:
                    try:
                        totals[col] = sum(r.get(col, 0) or 0 for r in all_data if isinstance(r.get(col), (int, float)))
                    except Exception:
                        totals[col] = 0

        return {
            "success": True,
            "data": data,
            "total": total_count,
            "page": request.page,
            "page_size": grid_page_size,
            "total_pages": (total_count + grid_page_size - 1) // grid_page_size,
            "totals": totals,
            "columns": columns,
        }
    except HTTPException:
        raise
    except Exception as e:
        return {"success": False, "error": str(e), "data": []}


@router.post("/grids/{grid_id}/export")
async def export_grid_data(
    grid_id: int,
    format: str = "csv",
    dwh_code:    Optional[str] = Header(None, alias="X-DWH-Code"),
    user_id_hdr: Optional[str] = Header(None, alias="X-User-Id"),
):
    """Exporte les donnees d'une grille

    Utilise le resolver avec fallback Template -> Override
    """
    from ..services.parameter_resolver import inject_params
    from ..services.permissions import enforce_report_access

    try:
        uid = int(user_id_hdr) if user_id_hdr else None
        enforce_report_access(uid, 'gridview', grid_id, dwh_code)
    except Exception as perm_err:
        if hasattr(perm_err, 'status_code'):
            raise

    try:
        # Recuperer la config (base client si dwh_code, sinon centrale)
        results = _gv_read(
            "SELECT * FROM APP_GridViews WHERE id = ?",
            (grid_id,),
            dwh_code=dwh_code,
        )
        if not results:
            raise HTTPException(status_code=404, detail="Grille non trouvee")

        grid = results[0]
        datasource_id = grid.get('data_source_id')
        datasource_code = grid.get('data_source_code')
        origin = None

        # Recuperer la source avec fallback Template -> Override
        try:
            if datasource_code:
                datasource = datasource_resolver.resolve_by_code(datasource_code, dwh_code)
            elif datasource_id:
                datasource = datasource_resolver.resolve_by_id(datasource_id, dwh_code)
            else:
                raise HTTPException(status_code=404, detail="Source de donnees non configuree")
            query = datasource.query_template
            origin = datasource.origin.value
        except ValueError:
            # Fallback: ancien comportement
            source = execute_query(
                "SELECT * FROM APP_DataSources WHERE id = ?",
                (datasource_id,),
                use_cache=False
            )
            if not source:
                raise HTTPException(status_code=404, detail="Source non trouvee")
            query = source[0]['query_template']

        # Injecter les params par defaut
        query = inject_params(query, {})

        # Router vers le bon DWH si template
        effective_dwh = dwh_code
        if not effective_dwh:
            try:
                dwh_list = execute_master_query(
                    "SELECT TOP 1 code FROM APP_DWH WHERE actif = 1 ORDER BY id",
                    use_cache=True
                )
                if dwh_list:
                    effective_dwh = dwh_list[0]['code']
            except Exception:
                pass

        # Executer la requete complete
        if effective_dwh:
            data = DWHConnectionManager.execute_dwh_query(effective_dwh, query, use_cache=False)
        else:
            data = execute_query(query, use_cache=False)

        return {
            "success": True,
            "data": data,
            "total": len(data),
            "format": format
        }
    except HTTPException:
        raise
    except Exception as e:
        return {"success": False, "error": str(e)}


# =============================================================================
# EXPORT POWERPOINT GRILLE
# =============================================================================

@router.get("/grids/{grid_id}/export/pptx")
async def export_grid_pptx(
    grid_id: int,
    dwh_code:    Optional[str] = Header(None, alias="X-DWH-Code"),
    user_id_hdr: Optional[str] = Header(None, alias="X-User-Id"),
):
    """Génère un fichier PowerPoint à partir des données de la grille (chargées côté serveur)."""
    from fastapi.responses import StreamingResponse
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from datetime import datetime
    from ..services.parameter_resolver import inject_params
    import io

    BLUE       = RGBColor(0x25, 0x63, 0xEB)
    DARK_BLUE  = RGBColor(0x1E, 0x40, 0xAF)
    WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
    GRAY_LIGHT = RGBColor(0xF3, 0xF4, 0xF6)

    try:
        # ── Récupérer config grille (base client si dwh_code, sinon centrale) ──
        grid_rows = _gv_read("SELECT * FROM APP_GridViews WHERE id = ?", (grid_id,), dwh_code=dwh_code)
        if not grid_rows:
            raise HTTPException(status_code=404, detail="Grille non trouvée")

        grid = grid_rows[0]
        grid_title = grid.get("nom") or "Grille"
        datasource_id   = grid.get("data_source_id")
        datasource_code = grid.get("data_source_code")

        # ── Colonnes depuis la config ──
        raw_cols = json.loads(grid.get("columns_config") or "[]")
        columns = [
            {"field": c.get("field", ""), "headerName": c.get("header") or c.get("headerName") or c.get("field", "")}
            for c in raw_cols
            if c.get("visible", True) and c.get("field")
        ]

        # ── Résoudre la source de données ──
        origin = None
        try:
            if datasource_code:
                datasource = datasource_resolver.resolve_by_code(datasource_code, dwh_code)
            elif datasource_id:
                datasource = datasource_resolver.resolve_by_id(datasource_id, dwh_code)
            else:
                raise HTTPException(status_code=404, detail="Source de données non configurée")
            query = datasource.query_template
            origin = datasource.origin.value
        except ValueError:
            source = execute_query("SELECT * FROM APP_DataSources WHERE id = ?", (datasource_id,), use_cache=False)
            if not source:
                raise HTTPException(status_code=404, detail="Source non trouvée")
            query = source[0]["query_template"]

        query = inject_params(query, {})

        effective_dwh = dwh_code
        if not effective_dwh:
            try:
                dwh_list = execute_master_query("SELECT TOP 1 code FROM APP_DWH WHERE actif = 1 ORDER BY id", use_cache=True)
                if dwh_list:
                    effective_dwh = dwh_list[0]["code"]
            except Exception:
                pass

        if effective_dwh:
            rows = DWHConnectionManager.execute_dwh_query(effective_dwh, query, use_cache=False)
        else:
            rows = execute_query(query, use_cache=False)

        # Si pas de colonnes définies, les déduire des données
        if not columns and rows:
            columns = [{"field": k, "headerName": k} for k in rows[0].keys()]

        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        # ── Slide titre ──
        slide = prs.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = DARK_BLUE

        txb = slide.shapes.add_textbox(Inches(1), Inches(2.6), Inches(11.33), Inches(1.2))
        p = txb.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = grid_title
        run.font.size = Pt(36); run.font.bold = True; run.font.color.rgb = WHITE

        txb2 = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(11.33), Inches(0.5))
        p2 = txb2.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = f"Généré le {datetime.now().strftime('%d/%m/%Y à %H:%M')}  —  OptiBoard"
        r2.font.size = Pt(14); r2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)

        if not columns or not rows:
            output = io.BytesIO()
            prs.save(output); output.seek(0)
            return StreamingResponse(output,
                media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                headers={"Content-Disposition": f'attachment; filename="grille_{grid_id}.pptx"'})

        # ── Slides données ──
        ROWS_PER_SLIDE = 22
        num_cols = len(columns)
        headers  = [c.get("headerName") or c.get("field", "") for c in columns]
        fields   = [c.get("field", "") for c in columns]

        chunks = [rows[i:i+ROWS_PER_SLIDE] for i in range(0, max(len(rows),1), ROWS_PER_SLIDE)]
        if not chunks:
            chunks = [[]]

        total_slides = len(chunks)
        for si, chunk in enumerate(chunks):
            slide = prs.slides.add_slide(blank)

            hdr = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(9), Inches(0.4))
            rh = hdr.text_frame.paragraphs[0].add_run()
            rh.text = grid_title; rh.font.size = Pt(11); rh.font.bold = True; rh.font.color.rgb = DARK_BLUE

            if total_slides > 1:
                pag = slide.shapes.add_textbox(Inches(10), Inches(0.15), Inches(3), Inches(0.4))
                pg = pag.text_frame.paragraphs[0]
                pg.alignment = PP_ALIGN.RIGHT
                rp = pg.add_run(); rp.text = f"{si+1} / {total_slides}"
                rp.font.size = Pt(9); rp.font.color.rgb = RGBColor(0x6B,0x72,0x80)

            num_rows_table = len(chunk) + 1
            if num_rows_table < 2:
                continue

            table = slide.shapes.add_table(
                num_rows_table, num_cols,
                Inches(0.3), Inches(0.65), Inches(12.73), Inches(6.6)
            ).table

            col_w = Inches(12.73) // num_cols
            for ci in range(num_cols):
                table.columns[ci].width = col_w

            for ci, h in enumerate(headers):
                cell = table.cell(0, ci)
                cell.fill.solid(); cell.fill.fore_color.rgb = BLUE
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                r = p.add_run(); r.text = str(h)
                r.font.size = Pt(8); r.font.bold = True; r.font.color.rgb = WHITE

            for ri, row_data in enumerate(chunk):
                for ci, field in enumerate(fields):
                    val = row_data.get(field, "")
                    cell = table.cell(ri + 1, ci)
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = GRAY_LIGHT if ri % 2 == 0 else WHITE
                    p = cell.text_frame.paragraphs[0]
                    is_num = isinstance(val, (int, float)) and not isinstance(val, bool)
                    p.alignment = PP_ALIGN.RIGHT if is_num else PP_ALIGN.LEFT
                    r = p.add_run()
                    r.text = f"{val:,.2f}" if is_num else str(val if val is not None else "")
                    r.font.size = Pt(8)

        output = io.BytesIO()
        prs.save(output); output.seek(0)
        return StreamingResponse(
            output,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="grille_{grid_id}.pptx"'}
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


# =============================================================================
# PREFERENCES UTILISATEUR PAR GRILLE
# =============================================================================

@router.get("/grids/{grid_id}/user-prefs/{user_id}")
async def get_user_prefs(grid_id: int, user_id: int):
    """Recupere les preferences de colonnes d'un utilisateur pour une grille"""
    try:
        results = execute_query(
            "SELECT columns_config FROM APP_GridView_User_Prefs WHERE grid_id = ? AND user_id = ?",
            (grid_id, user_id),
            use_cache=False
        )
        if results:
            columns = json.loads(results[0].get('columns_config') or '[]')
            return {"success": True, "data": columns, "has_prefs": True}
        return {"success": True, "data": [], "has_prefs": False}
    except Exception as e:
        return {"success": False, "error": str(e), "data": [], "has_prefs": False}


@router.put("/grids/{grid_id}/user-prefs/{user_id}")
async def save_user_prefs(grid_id: int, user_id: int, body: Dict[str, Any]):
    """Sauvegarde les preferences de colonnes d'un utilisateur pour une grille"""
    try:
        columns = body.get('columns', [])
        columns_json = json.dumps(columns, ensure_ascii=False)

        # UPSERT : update si existe, insert sinon
        existing = execute_query(
            "SELECT id FROM APP_GridView_User_Prefs WHERE grid_id = ? AND user_id = ?",
            (grid_id, user_id),
            use_cache=False
        )

        with get_db_cursor() as cursor:
            if existing:
                cursor.execute(
                    "UPDATE APP_GridView_User_Prefs SET columns_config = ?, updated_at = GETDATE() WHERE grid_id = ? AND user_id = ?",
                    (columns_json, grid_id, user_id)
                )
            else:
                cursor.execute(
                    "INSERT INTO APP_GridView_User_Prefs (grid_id, user_id, columns_config) VALUES (?, ?, ?)",
                    (grid_id, user_id, columns_json)
                )

        return {"success": True, "message": "Preferences sauvegardees"}
    except Exception as e:
        print(f"[ERROR] save_user_prefs: {e}")
        return {"success": False, "error": str(e)}


@router.delete("/grids/{grid_id}/user-prefs/{user_id}")
async def reset_user_prefs(grid_id: int, user_id: int):
    """Reinitialise les preferences utilisateur (retour a la config par defaut)"""
    try:
        with get_db_cursor() as cursor:
            cursor.execute(
                "DELETE FROM APP_GridView_User_Prefs WHERE grid_id = ? AND user_id = ?",
                (grid_id, user_id)
            )
        return {"success": True, "message": "Preferences reinitialisees"}
    except Exception as e:
        return {"success": False, "error": str(e)}
